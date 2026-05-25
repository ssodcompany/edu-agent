"""
BNI SoR v_final 9.1차 — native PPTX 빌더 (4명 평가단 fix 반영)
========================================

페이퍼로즈 SOP 준수:
- python-pptx로 native shape·텍스트 직접 생성 (이미지 박는 거 아님)
- Pretendard 폰트명 명시 (사용자 시스템 자동 로드)
- Claude DESIGN.md 토큰 (canvas #faf9f5 · coral #cc785c · ink #141413)
- 16:9 (13.333 × 7.5 inch)
- 슬라이드 27장 (v9 28장 → v9.1 27장, s26+s27 통합)

v9.1 fix:
- F1: s26·s27 통합 → 단일 듀얼 모델 슬라이드 (inclusion 톤)
- F2: s20 팔란티어 시총·자기 등급 자랑 톤 제거
- F3: s21 비교표 4행 차별화 ("핵심" 행 동일 → 한국 OS 맥락)
- F4: coral 절제 — 강조 슬라이드 7개만 (s14·s17·s19·s20·s25·s26·s27)
- F5: referral hook 단일화

산출물: BNI_SoR_v_final_9_1.pptx (편집 가능 native)
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ============================================================
# 토큰 (Claude DESIGN.md)
# ============================================================
C_CANVAS       = RGBColor(0xFA, 0xF9, 0xF5)  # #faf9f5
C_SURFACE_SOFT = RGBColor(0xF5, 0xF0, 0xE8)  # #f5f0e8
C_SURFACE_CARD = RGBColor(0xEF, 0xE9, 0xDE)  # #efe9de
C_SURFACE_DARK = RGBColor(0x18, 0x17, 0x15)  # #181715
C_SURFACE_DARK_ELEVATED = RGBColor(0x25, 0x23, 0x20)  # #252320
C_INK          = RGBColor(0x14, 0x14, 0x13)  # #141413
C_BODY         = RGBColor(0x3D, 0x3D, 0x3A)  # #3d3d3a
C_BODY_STRONG  = RGBColor(0x25, 0x25, 0x23)  # #252523
C_MUTED        = RGBColor(0x6C, 0x6A, 0x64)  # #6c6a64
C_MUTED_SOFT   = RGBColor(0x7A, 0x77, 0x6F)  # #7a776f (WCAG AA tuned)
C_CORAL        = RGBColor(0xCC, 0x78, 0x5C)  # #cc785c
C_CORAL_ACTIVE = RGBColor(0xA9, 0x58, 0x3E)  # #a9583e
C_HAIRLINE     = RGBColor(0xE6, 0xDF, 0xD8)  # #e6dfd8
C_ON_PRIMARY   = RGBColor(0xFF, 0xFF, 0xFF)  # white
C_ON_DARK      = RGBColor(0xFA, 0xF9, 0xF5)  # cream-on-dark
C_ON_DARK_SOFT = RGBColor(0xA0, 0x9D, 0x96)  # #a09d96

FONT_KR   = "Pretendard"
FONT_KR_LIGHT = "Pretendard"  # Variable; weight controlled via bold
FONT_MONO = "SF Mono"
FONT_FALLBACK_KR = "Apple SD Gothic Neo"

# ============================================================
# 슬라이드 크기 (16:9)
# ============================================================
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# 안전영역 (4% margin)
SAFE_X = Inches(0.667)
SAFE_Y = Inches(0.5)
CONTENT_W = SLIDE_W - SAFE_X * 2
CONTENT_H = SLIDE_H - SAFE_Y * 2

# Asset 폴더
ASSETS_DIR = Path(__file__).parent / "assets"
OUT_PATH = Path(__file__).parent / "BNI_SoR_v_final_9_1.pptx"
TOTAL_SLIDES_STR = "27"


# ============================================================
# Helpers
# ============================================================

def set_fill(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color

def no_fill(shape):
    shape.fill.background()

def set_line(shape, color=None, width_pt: float = 0):
    line = shape.line
    if width_pt == 0:
        line.fill.background()
    else:
        line.color.rgb = color
        line.width = Pt(width_pt)

def set_text(
    text_frame,
    text: str,
    *,
    font=FONT_KR,
    size: int = 14,
    color: RGBColor = C_INK,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    spacing: float = 1.0,  # line height (NOT line spacing pt)
    letter_spacing: int = 0,  # in 1/100 pt, negative = tighter
):
    text_frame.word_wrap = True
    text_frame.margin_left = Pt(0)
    text_frame.margin_right = Pt(0)
    text_frame.margin_top = Pt(0)
    text_frame.margin_bottom = Pt(0)
    # auto_size 비활성 (수동 컨트롤)
    try:
        from pptx.enum.text import MSO_AUTO_SIZE
        text_frame.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        pass

    p = text_frame.paragraphs[0]
    p.alignment = align
    p.line_spacing = spacing
    # clear runs
    for run in list(p.runs):
        run.text = ""
    run = p.add_run() if not p.runs else p.runs[0]
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    # East Asian font fallback (XML hack)
    rPr = run._r.get_or_add_rPr()
    latin = rPr.find(qn("a:latin"))
    if latin is None:
        latin = etree.SubElement(rPr, qn("a:latin"))
    latin.set("typeface", font)
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = etree.SubElement(rPr, qn("a:ea"))
    ea.set("typeface", font)
    cs = rPr.find(qn("a:cs"))
    if cs is None:
        cs = etree.SubElement(rPr, qn("a:cs"))
    cs.set("typeface", font)
    # Letter spacing (spc attribute, units = 1/100 pt)
    if letter_spacing != 0:
        rPr.set("spc", str(letter_spacing))


def add_textbox(
    slide,
    x, y, w, h,
    text: str,
    *,
    font=FONT_KR,
    size: int = 14,
    color: RGBColor = C_INK,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    spacing: float = 1.3,
    letter_spacing: int = 0,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.vertical_anchor = anchor
    set_text(
        tf, text,
        font=font, size=size, color=color, bold=bold,
        align=align, spacing=spacing, letter_spacing=letter_spacing,
    )
    return box


def add_multiline_textbox(
    slide,
    x, y, w, h,
    lines,  # list[ dict(text, font, size, color, bold, align) ]
    *,
    anchor=MSO_ANCHOR.TOP,
    base_spacing: float = 1.4,
):
    """여러 줄을 paragraph로 쌓는 textbox."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.vertical_anchor = anchor
    tf.word_wrap = True
    tf.margin_left = Pt(0); tf.margin_right = Pt(0)
    tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
    try:
        from pptx.enum.text import MSO_AUTO_SIZE
        tf.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        pass

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = line.get("align", PP_ALIGN.LEFT)
        p.line_spacing = line.get("spacing", base_spacing)
        if line.get("space_before"):
            p.space_before = Pt(line["space_before"])
        if line.get("space_after"):
            p.space_after = Pt(line["space_after"])
        # clear existing runs
        for run in list(p.runs):
            run.text = ""
        run = p.add_run() if not p.runs else p.runs[0]
        run.text = line["text"]
        run.font.name = line.get("font", FONT_KR)
        run.font.size = Pt(line.get("size", 14))
        run.font.color.rgb = line.get("color", C_INK)
        run.font.bold = line.get("bold", False)
        rPr = run._r.get_or_add_rPr()
        for tag in ("a:latin", "a:ea", "a:cs"):
            el = rPr.find(qn(tag))
            if el is None:
                el = etree.SubElement(rPr, qn(tag))
            el.set("typeface", line.get("font", FONT_KR))
        if "letter_spacing" in line and line["letter_spacing"] != 0:
            rPr.set("spc", str(line["letter_spacing"]))
    return box


def add_rect(slide, x, y, w, h, *, fill=None, line_color=None, line_width=0, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, x, y, w, h)
    if radius is not None:
        # ROUNDED_RECTANGLE adjustment: smaller value = smaller radius. 0~1 ratio of half-min-dim.
        # We approximate: radius_emu / min(w,h)/2
        try:
            target = float(radius) / float(min(w, h) / 2)
            target = max(0.0, min(1.0, target))
            shape.adjustments[0] = target
        except Exception:
            pass
    if fill is None:
        no_fill(shape)
    else:
        set_fill(shape, fill)
    if line_color is None or line_width == 0:
        set_line(shape)
    else:
        set_line(shape, line_color, line_width)
    return shape


def add_line(slide, x1, y1, x2, y2, *, color=C_HAIRLINE, width_pt=1.0):
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)  # 1 = STRAIGHT
    connector.line.color.rgb = color
    connector.line.width = Pt(width_pt)
    return connector


def add_canvas_background(slide, color=C_CANVAS):
    bg = add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=color)
    return bg


def add_dark_background(slide):
    return add_canvas_background(slide, C_SURFACE_DARK)


def add_brandlogy(slide, *, page_no: str, total: str = TOTAL_SLIDES_STR, dark: bool = False, eyebrow: str = "", brand_mark: bool = True, source_line: str = ""):
    """4-corner brandlogy 패턴 — 좌상 eyebrow / 우상 brand-mark / 좌하 page-no / 우하 source."""
    text_color = C_ON_DARK_SOFT if dark else C_MUTED_SOFT
    # 좌상 eyebrow
    if eyebrow:
        add_textbox(
            slide,
            SAFE_X, SAFE_Y,
            Inches(6), Pt(20),
            eyebrow,
            font=FONT_MONO, size=11, color=text_color,
            letter_spacing=150,
        )
    # 우상 brand-mark (SSOD)
    if brand_mark:
        bm_w = Inches(1.2)
        bm_h = Pt(20)
        add_textbox(
            slide,
            SLIDE_W - SAFE_X - bm_w, SAFE_Y,
            bm_w, bm_h,
            "SSOD",
            font=FONT_KR, size=12, color=text_color, bold=True,
            align=PP_ALIGN.RIGHT,
            letter_spacing=200,
        )
    # 좌하 page-no
    add_textbox(
        slide,
        SAFE_X, SLIDE_H - SAFE_Y - Pt(20),
        Inches(2), Pt(20),
        f"{page_no} / {total}",
        font=FONT_MONO, size=10, color=text_color,
    )
    # 우하 source
    if source_line:
        add_textbox(
            slide,
            SLIDE_W - SAFE_X - Inches(6), SLIDE_H - SAFE_Y - Pt(20),
            Inches(6), Pt(20),
            source_line,
            font=FONT_KR, size=10, color=text_color,
            align=PP_ALIGN.RIGHT,
        )


# ============================================================
# 슬라이드 빌더
# ============================================================

def new_slide(prs, *, dark=False):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    if dark:
        add_dark_background(slide)
    else:
        add_canvas_background(slide)
    return slide


# --- 1. cover (s1)
def slide_cover(prs, page_no):
    slide = new_slide(prs, dark=False)

    # 좌상 spike-mark + wordmark (텍스트로 대체)
    add_textbox(
        slide,
        SAFE_X, SAFE_Y,
        Inches(4), Pt(22),
        "■  SSOD COMPANY",
        font=FONT_KR, size=12, color=C_INK, bold=True,
        letter_spacing=300,
    )

    # 중앙 large display (Copernicus 대체 — Pretendard tight)
    title_y = Inches(2.2)
    add_textbox(
        slide,
        SAFE_X, title_y,
        CONTENT_W, Inches(1.6),
        "구조가 없으면,",
        font=FONT_KR, size=72, color=C_INK, bold=False,
        align=PP_ALIGN.LEFT,
        letter_spacing=-350,
    )
    add_textbox(
        slide,
        SAFE_X, title_y + Inches(1.4),
        CONTENT_W, Inches(1.6),
        "AI도 없습니다.",
        font=FONT_KR, size=72, color=C_CORAL, bold=False,
        align=PP_ALIGN.LEFT,
        letter_spacing=-350,
    )

    # 하단 subtitle
    add_multiline_textbox(
        slide,
        SAFE_X, Inches(5.8),
        CONTENT_W, Inches(1),
        [
            {"text": "회사의 일하는 구조 — 데이터 · 일 · 판단", "font": FONT_KR, "size": 16, "color": C_BODY},
            {"text": "BNI K-Chapter Feature Presentation · 김지명", "font": FONT_MONO, "size": 11, "color": C_MUTED_SOFT, "letter_spacing": 150, "space_before": 6},
        ]
    )

    add_brandlogy(slide, page_no=page_no, eyebrow="01 · COVER", brand_mark=False)
    return slide


# --- 2. intro centered (s2)
def slide_intro_centered(prs, page_no):
    slide = new_slide(prs)
    add_brandlogy(slide, page_no=page_no, eyebrow="02 · 자기소개")

    # 큰 이름
    add_textbox(
        slide, SAFE_X, Inches(1.6), CONTENT_W, Inches(1.4),
        "김지명",
        font=FONT_KR, size=96, color=C_INK, bold=False,
        align=PP_ALIGN.CENTER, letter_spacing=-300,
    )
    # 소속 · 직책
    add_textbox(
        slide, SAFE_X, Inches(3.0), CONTENT_W, Inches(0.4),
        "쏘드 컴퍼니 · 대표",
        font=FONT_KR, size=18, color=C_MUTED,
        align=PP_ALIGN.CENTER, letter_spacing=200,
    )

    # B2B 7년차 pill
    pill_w = Inches(2.5)
    pill_h = Inches(0.45)
    pill = add_rect(
        slide,
        (SLIDE_W - pill_w) / 2, Inches(3.55),
        pill_w, pill_h,
        fill=C_SURFACE_CARD, radius=999,
    )
    add_textbox(
        slide, (SLIDE_W - pill_w) / 2, Inches(3.6),
        pill_w, Inches(0.35),
        "B2B · 7년차 · 상주 컨설팅",
        font=FONT_KR, size=13, color=C_INK, bold=True,
        align=PP_ALIGN.CENTER, letter_spacing=150,
    )

    # divider
    add_line(slide, Inches(6.2), Inches(4.5), Inches(7.1), Inches(4.5), color=C_HAIRLINE, width_pt=2)

    # eyebrow + headline + sub
    add_textbox(
        slide, SAFE_X, Inches(4.85), CONTENT_W, Inches(0.3),
        "하는 일",
        font=FONT_MONO, size=11, color=C_MUTED_SOFT,
        align=PP_ALIGN.CENTER, letter_spacing=200,
    )
    add_textbox(
        slide, SAFE_X, Inches(5.2), CONTENT_W, Inches(0.7),
        "회사의 일하는 구조를 만듭니다.",
        font=FONT_KR, size=40, color=C_CORAL, bold=False,
        align=PP_ALIGN.CENTER, letter_spacing=-150,
    )
    add_textbox(
        slide, SAFE_X, Inches(6.05), CONTENT_W, Inches(0.4),
        "데이터 · 일 · 판단 — 셋이 한 자리에 맞물리게.",
        font=FONT_KR, size=17, color=C_BODY,
        align=PP_ALIGN.CENTER, letter_spacing=-50,
    )
    return slide


# --- 3. logo wall (s3)
def slide_logo_wall(prs, page_no):
    slide = new_slide(prs)
    add_brandlogy(slide, page_no=page_no, eyebrow="03 · 함께한 자리들")

    # 헤더
    add_textbox(
        slide, SAFE_X, Inches(0.95), CONTENT_W, Inches(0.5),
        "도구 · 협업 · 강연",
        font=FONT_KR, size=14, color=C_MUTED, align=PP_ALIGN.LEFT,
        letter_spacing=200,
    )
    add_textbox(
        slide, SAFE_X, Inches(1.4), CONTENT_W, Inches(0.7),
        "같은 진단 렌즈로 — 다양한 규모·직군에서.",
        font=FONT_KR, size=28, color=C_INK, bold=False,
        align=PP_ALIGN.LEFT, letter_spacing=-150,
    )

    # 4 × 2 grid
    cards = [
        # row 1: 도구
        ("Notion",         "일하는 자리",        "공식 앰버서더"),
        ("n8n",            "자동화 혈관",        "공식 앰버서더"),
        ("Claude",         "판단 시스템",        "Anthropic"),
        ("Microsoft Excel", "데이터 인입",       "Office 365"),
        # row 2: 협업·강연
        ("Notion Camp",     "HR 본사 워크숍",     "2025"),
        ("상장사 onsite",   "6개월 상주 컨설팅",  "NDA · 익명"),
        ("Notion·n8n 공식 행사", "강연",         "Seoul · Online"),
        ("자유 슬롯",       "발표 현장 photo",    "사용자 drop"),
    ]
    grid_x = SAFE_X
    grid_y = Inches(2.5)
    grid_w = CONTENT_W
    cols = 4
    rows = 2
    gap = Inches(0.18)
    cell_w = (grid_w - gap * (cols - 1)) / cols
    cell_h = Inches(2.0)

    for idx, (name, role, sub) in enumerate(cards):
        r = idx // cols
        c = idx % cols
        x = grid_x + (cell_w + gap) * c
        y = grid_y + (cell_h + gap) * r
        is_placeholder = idx >= 7
        is_partner = idx in (4, 5, 6)
        if is_placeholder:
            card = add_rect(slide, x, y, cell_w, cell_h, fill=C_CANVAS, line_color=C_HAIRLINE, line_width=1.5, radius=Inches(0.15))
        else:
            card = add_rect(slide, x, y, cell_w, cell_h, fill=C_SURFACE_CARD, radius=Inches(0.15))

        # name
        add_textbox(
            slide, x, y + Inches(0.4), cell_w, Inches(0.5),
            name,
            font=FONT_KR, size=22, color=C_INK if not is_placeholder else C_MUTED_SOFT, bold=True,
            align=PP_ALIGN.CENTER, letter_spacing=-100,
        )
        # role
        add_textbox(
            slide, x, y + Inches(1.0), cell_w, Inches(0.35),
            role,
            font=FONT_KR, size=13, color=C_BODY if not is_placeholder else C_MUTED_SOFT,
            align=PP_ALIGN.CENTER,
        )
        # sub mono
        add_textbox(
            slide, x, y + Inches(1.4), cell_w, Inches(0.3),
            sub,
            font=FONT_MONO, size=10, color=C_MUTED_SOFT,
            align=PP_ALIGN.CENTER, letter_spacing=150,
        )
    return slide


# --- 4. ambassador PUNCH (s4)
def slide_ambassador_punch(prs, page_no):
    slide = new_slide(prs)
    add_brandlogy(slide, page_no=page_no, eyebrow="04 · 자격")

    # eyebrow
    add_textbox(
        slide, SAFE_X, Inches(1.2), CONTENT_W, Inches(0.4),
        "GLOBAL · OFFICIAL",
        font=FONT_MONO, size=12, color=C_CORAL, bold=True,
        align=PP_ALIGN.CENTER, letter_spacing=350,
    )

    # 큰 punch
    add_multiline_textbox(
        slide,
        SAFE_X, Inches(1.9),
        CONTENT_W, Inches(3.0),
        [
            {"text": "둘 다 — 글로벌 공식 앰버서더.", "font": FONT_KR, "size": 56, "color": C_INK, "align": PP_ALIGN.CENTER, "letter_spacing": -250, "spacing": 1.1},
            {"text": "한국에서 한 명.", "font": FONT_KR, "size": 56, "color": C_CORAL, "bold": False, "align": PP_ALIGN.CENTER, "letter_spacing": -250, "spacing": 1.1, "space_before": 8},
        ]
    )

    # 두 배지 row (텍스트 카드로 표현)
    badge_w = Inches(4.0)
    badge_h = Inches(1.0)
    badge_gap = Inches(0.4)
    badges_total_w = badge_w * 2 + badge_gap
    badges_x = (SLIDE_W - badges_total_w) / 2

    # Notion badge
    add_rect(slide, badges_x, Inches(5.5), badge_w, badge_h, fill=C_INK, radius=Inches(0.1))
    add_textbox(
        slide, badges_x, Inches(5.7), badge_w, Inches(0.3),
        "NOTION",
        font=FONT_KR, size=14, color=C_ON_DARK, bold=True,
        align=PP_ALIGN.CENTER, letter_spacing=400,
    )
    add_textbox(
        slide, badges_x, Inches(6.05), badge_w, Inches(0.3),
        "Official Ambassador",
        font=FONT_KR, size=12, color=C_ON_DARK_SOFT,
        align=PP_ALIGN.CENTER, letter_spacing=200,
    )

    # n8n badge
    n8n_x = badges_x + badge_w + badge_gap
    add_rect(slide, n8n_x, Inches(5.5), badge_w, badge_h, fill=C_INK, radius=Inches(0.1))
    add_textbox(
        slide, n8n_x, Inches(5.7), badge_w, Inches(0.3),
        "n8n",
        font=FONT_KR, size=14, color=C_ON_DARK, bold=True,
        align=PP_ALIGN.CENTER, letter_spacing=400,
    )
    add_textbox(
        slide, n8n_x, Inches(6.05), badge_w, Inches(0.3),
        "Official Ambassador",
        font=FONT_KR, size=12, color=C_ON_DARK_SOFT,
        align=PP_ALIGN.CENTER, letter_spacing=200,
    )
    return slide


# --- 5. 활동·도구 (s5)
def slide_activity(prs, page_no):
    slide = new_slide(prs)
    add_brandlogy(slide, page_no=page_no, eyebrow="05 · 활동")

    add_textbox(
        slide, SAFE_X, Inches(1.2), CONTENT_W, Inches(0.4),
        "도구 · 세 가지",
        font=FONT_MONO, size=12, color=C_MUTED, letter_spacing=200,
    )
    add_textbox(
        slide, SAFE_X, Inches(1.7), CONTENT_W, Inches(1.0),
        "회사가 어떻게 일해야 하는지 — 시스템으로 박아 둡니다.",
        font=FONT_KR, size=32, color=C_INK, letter_spacing=-150,
    )

    # 3-card grid: Notion · n8n · Claude Code
    tools = [
        ("Notion",       "일하는 자리",  "메시지·결재·문서 한 자리"),
        ("n8n",          "혈관",         "데이터 자동 흐름"),
        ("Claude Code",  "판단 시스템",  "사람 머릿속 룰 → 코드"),
    ]
    grid_y = Inches(3.6)
    cols = 3
    gap = Inches(0.3)
    cell_w = (CONTENT_W - gap * (cols - 1)) / cols
    cell_h = Inches(2.6)

    for i, (name, role, body) in enumerate(tools):
        x = SAFE_X + (cell_w + gap) * i
        add_rect(slide, x, grid_y, cell_w, cell_h, fill=C_SURFACE_CARD, radius=Inches(0.18))
        # mono frame
        add_textbox(
            slide, x + Inches(0.3), grid_y + Inches(0.3), cell_w - Inches(0.6), Inches(0.3),
            f"0{i+1}",
            font=FONT_MONO, size=11, color=C_CORAL, bold=True, letter_spacing=200,
        )
        # name
        add_textbox(
            slide, x + Inches(0.3), grid_y + Inches(0.7), cell_w - Inches(0.6), Inches(0.6),
            name,
            font=FONT_KR, size=32, color=C_INK, bold=True, letter_spacing=-150,
        )
        # role
        add_textbox(
            slide, x + Inches(0.3), grid_y + Inches(1.5), cell_w - Inches(0.6), Inches(0.4),
            role,
            font=FONT_KR, size=15, color=C_CORAL, bold=True, letter_spacing=200,
        )
        # body
        add_textbox(
            slide, x + Inches(0.3), grid_y + Inches(1.95), cell_w - Inches(0.6), Inches(0.6),
            body,
            font=FONT_KR, size=14, color=C_BODY, spacing=1.4,
        )

    # 하단 thesis
    add_textbox(
        slide, SAFE_X, Inches(6.5), CONTENT_W, Inches(0.4),
        "→ 회사가 어떻게 일해야 하는지를 시스템으로.",
        font=FONT_KR, size=16, color=C_MUTED, letter_spacing=-50,
    )
    return slide


# --- 6. opening hook (s6)
def slide_opening_hook(prs, page_no):
    slide = new_slide(prs)
    add_brandlogy(slide, page_no=page_no, eyebrow="06 · 오프닝 훅")

    add_textbox(
        slide, SAFE_X, Inches(0.95), CONTENT_W, Inches(0.4),
        "세 가지만 묻고 시작하겠습니다.",
        font=FONT_KR, size=15, color=C_MUTED, letter_spacing=-50,
    )

    questions = [
        ("Q1.", "어제 결정한 일 — 어디 적혀 있나요?"),
        ("Q2.", "고객 견적 요청 — 어디로 갑니까?"),
        ("Q3.", "지난달 매출 1등이 누구한테서 왔는지 — 한 번에 답 가능하세요?"),
    ]
    qy = Inches(1.7)
    qh = Inches(1.3)
    qgap = Inches(0.25)

    for i, (qno, q) in enumerate(questions):
        y = qy + (qh + qgap) * i
        # mono left
        add_textbox(
            slide, SAFE_X, y + Inches(0.3), Inches(1.2), Inches(0.6),
            qno,
            font=FONT_MONO, size=28, color=C_CORAL, bold=True, letter_spacing=200,
        )
        # body
        add_textbox(
            slide, SAFE_X + Inches(1.3), y + Inches(0.15), CONTENT_W - Inches(1.3), Inches(1.0),
            q,
            font=FONT_KR, size=28, color=C_INK, letter_spacing=-150, spacing=1.3,
        )
        # hairline divider
        if i < len(questions) - 1:
            add_line(
                slide,
                SAFE_X, y + qh,
                SAFE_X + CONTENT_W, y + qh,
                color=C_HAIRLINE, width_pt=1,
            )

    # 클로징 cue
    add_textbox(
        slide, SAFE_X, Inches(6.0), CONTENT_W, Inches(0.6),
        "하나라도 \"음, 어디 있더라…\" 하셨다면 —",
        font=FONT_KR, size=18, color=C_MUTED, spacing=1.4,
    )
    add_textbox(
        slide, SAFE_X, Inches(6.4), CONTENT_W, Inches(0.6),
        "그 회사에 ChatGPT 깔아도, 일이 크게 바뀌진 않습니다.",
        font=FONT_KR, size=20, color=C_INK, bold=True, letter_spacing=-50,
    )
    return slide


# --- 7~12. SoR/SoE/SoI definition + insight 6장
def slide_so_def(prs, page_no, *, frame: str, frame_letter: str, no: str, name_kr: str, body: str, examples: str, is_climax_frame: bool = False):
    """SoR/SoE/SoI 정의 카드. is_climax_frame=True (s11 SoI 정의)에서만 coral 강조."""
    slide = new_slide(prs)
    add_brandlogy(slide, page_no=page_no, eyebrow=f"{no} · {name_kr} — 정의")

    # 좌측 frame letter 큰 라벨 (coral은 SoI thesis 진입 신호에만 — v9.1 F4 coral 절제)
    frame_color = C_CORAL if is_climax_frame else C_INK
    add_textbox(
        slide, SAFE_X, Inches(1.3), Inches(5.5), Inches(0.5),
        frame,
        font=FONT_MONO, size=56, color=frame_color, bold=True, letter_spacing=200, spacing=1.0,
    )
    add_textbox(
        slide, SAFE_X, Inches(2.5), Inches(5.5), Inches(0.5),
        f"System of {frame_letter}",
        font=FONT_MONO, size=14, color=C_MUTED_SOFT, letter_spacing=150,
    )

    # 우측 카드
    card_x = SAFE_X + Inches(6)
    card_y = Inches(1.3)
    card_w = CONTENT_W - Inches(6)
    card_h = Inches(5.3)
    add_rect(slide, card_x, card_y, card_w, card_h, fill=C_SURFACE_CARD, radius=Inches(0.2))

    # 카드 내부
    inner_x = card_x + Inches(0.6)
    inner_w = card_w - Inches(1.2)

    name_color = C_CORAL if is_climax_frame else C_INK
    add_textbox(
        slide, inner_x, card_y + Inches(0.7), inner_w, Inches(0.4),
        name_kr,
        font=FONT_KR, size=18, color=name_color, bold=True, letter_spacing=100,
    )
    add_textbox(
        slide, inner_x, card_y + Inches(1.2), inner_w, Inches(2.0),
        body,
        font=FONT_KR, size=22, color=C_INK, letter_spacing=-50, spacing=1.4,
    )
    # divider
    add_line(slide, inner_x, card_y + Inches(3.7), inner_x + inner_w, card_y + Inches(3.7), color=C_HAIRLINE, width_pt=1)

    add_textbox(
        slide, inner_x, card_y + Inches(3.95), inner_w, Inches(0.3),
        "예시",
        font=FONT_MONO, size=11, color=C_MUTED_SOFT, letter_spacing=200,
    )
    add_textbox(
        slide, inner_x, card_y + Inches(4.3), inner_w, Inches(0.9),
        examples,
        font=FONT_KR, size=15, color=C_BODY, spacing=1.4,
    )
    return slide


def slide_so_insight(prs, page_no, *, frame: str, no: str, headline: str, body: str, is_thesis: bool = False):
    """SoR/SoE/SoI 통찰 카드. is_thesis=True (s12 thesis)에서만 coral 강조."""
    slide = new_slide(prs)
    add_brandlogy(slide, page_no=page_no, eyebrow=f"{no} · 통찰")

    # frame letter mono (coral은 thesis 비트(s12)에만 — v9.1 F4)
    frame_color = C_CORAL if is_thesis else C_MUTED
    add_textbox(
        slide, SAFE_X, Inches(1.2), Inches(4), Inches(0.4),
        frame,
        font=FONT_MONO, size=22, color=frame_color, bold=True, letter_spacing=300,
    )

    # 큰 headline
    add_textbox(
        slide, SAFE_X, Inches(2.0), CONTENT_W, Inches(2.2),
        headline,
        font=FONT_KR, size=44, color=C_INK, letter_spacing=-200, spacing=1.2,
    )

    # 본문
    add_textbox(
        slide, SAFE_X, Inches(4.8), CONTENT_W, Inches(2.0),
        body,
        font=FONT_KR, size=18, color=C_BODY, spacing=1.5, letter_spacing=-50,
    )
    return slide


# --- 13. AI 시대 3축 연결 (s13)
def slide_3axis_connect(prs, page_no):
    slide = new_slide(prs)
    add_brandlogy(slide, page_no=page_no, eyebrow="13 · AI 시대 — 3축 연결")

    add_textbox(
        slide, SAFE_X, Inches(0.95), CONTENT_W, Inches(0.4),
        "AI 시대에 필요한 것은 — 셋, 그리고 그 연결.",
        font=FONT_KR, size=15, color=C_MUTED, letter_spacing=-50,
    )

    add_textbox(
        slide, SAFE_X, Inches(1.5), CONTENT_W, Inches(0.9),
        "셋을 갖추는 것 + 셋을 잇는 일.",
        font=FONT_KR, size=36, color=C_INK, letter_spacing=-200, spacing=1.2,
    )

    # 3 cards + connector
    grid_y = Inches(3.0)
    cols = 3
    gap = Inches(0.5)
    cell_w = (CONTENT_W - gap * (cols - 1)) / cols
    cell_h = Inches(2.4)

    axes = [
        ("SoR", "데이터",       "회사에 적힌 사실"),
        ("SoE", "일하는 자리",  "사람들이 모이는 공간"),
        ("SoI", "판단",         "결정이 내려지는 자리"),
    ]
    for i, (frame, role, body) in enumerate(axes):
        x = SAFE_X + (cell_w + gap) * i
        add_rect(slide, x, grid_y, cell_w, cell_h, fill=C_SURFACE_CARD, radius=Inches(0.2))
        # frame mono (v9.1 F4 — coral 절제, ink로 변경)
        add_textbox(
            slide, x, grid_y + Inches(0.3), cell_w, Inches(0.4),
            frame,
            font=FONT_MONO, size=28, color=C_INK, bold=True, align=PP_ALIGN.CENTER, letter_spacing=200,
        )
        # role
        add_textbox(
            slide, x, grid_y + Inches(0.9), cell_w, Inches(0.5),
            role,
            font=FONT_KR, size=24, color=C_INK, bold=True, align=PP_ALIGN.CENTER, letter_spacing=-100,
        )
        # body
        add_textbox(
            slide, x, grid_y + Inches(1.6), cell_w, Inches(0.5),
            body,
            font=FONT_KR, size=13, color=C_BODY, align=PP_ALIGN.CENTER,
        )

    # connector lines between 3 cards (v9.1 F4 — coral 절제, muted로)
    cy = grid_y + cell_h + Inches(0.3)
    add_line(slide, SAFE_X + cell_w * 0.5, cy, SAFE_X + cell_w * 1.5 + gap, cy, color=C_MUTED_SOFT, width_pt=2)
    add_line(slide, SAFE_X + cell_w * 1.5 + gap, cy, SAFE_X + cell_w * 2.5 + gap * 2, cy, color=C_MUTED_SOFT, width_pt=2)

    # 클로징 thesis
    add_textbox(
        slide, SAFE_X, Inches(6.4), CONTENT_W, Inches(0.6),
        "셋을 잇는 자리 — 그게 SSOD입니다.",
        font=FONT_KR, size=22, color=C_INK, bold=True, align=PP_ALIGN.CENTER, letter_spacing=-50,
    )
    return slide


# --- 14. SSOD reveal (s14)
def slide_ssod_reveal(prs, page_no):
    slide = new_slide(prs)
    add_brandlogy(slide, page_no=page_no, eyebrow="14 · SSOD")

    # eyebrow
    add_textbox(
        slide, SAFE_X, Inches(1.0), CONTENT_W, Inches(0.4),
        "쏘드 컴퍼니, SSOD라고도 씁니다.",
        font=FONT_KR, size=18, color=C_MUTED, align=PP_ALIGN.CENTER,
    )

    # 4-letter punch
    letters_y = Inches(2.0)
    letters_h = Inches(2.5)
    letter_w = CONTENT_W / 4
    letters = ["S", "S", "O", "D"]
    sublabels = ["Sync", "with", "Self · Other ·", "Data"]
    for i, (L, sub) in enumerate(zip(letters, sublabels)):
        x = SAFE_X + letter_w * i
        add_textbox(
            slide, x, letters_y, letter_w, Inches(1.6),
            L,
            font=FONT_KR, size=140, color=C_INK if i != 3 else C_CORAL, align=PP_ALIGN.CENTER,
            letter_spacing=-400, spacing=1.0,
        )
        add_textbox(
            slide, x, letters_y + Inches(1.7), letter_w, Inches(0.4),
            sub,
            font=FONT_MONO, size=14, color=C_MUTED_SOFT, align=PP_ALIGN.CENTER, letter_spacing=200,
        )

    # 하단 thesis
    add_textbox(
        slide, SAFE_X, Inches(5.5), CONTENT_W, Inches(0.5),
        "Sync with Self · Other · Data",
        font=FONT_MONO, size=18, color=C_INK, bold=True, align=PP_ALIGN.CENTER, letter_spacing=300,
    )
    add_textbox(
        slide, SAFE_X, Inches(6.1), CONTENT_W, Inches(0.5),
        "개인 · 조직 · 데이터 — 셋을 한 자리에 맞추는 일.",
        font=FONT_KR, size=18, color=C_BODY, align=PP_ALIGN.CENTER, letter_spacing=-50,
    )
    return slide


# --- 15. SSOD ↔ SoX 매핑 (s15)
def slide_ssod_mapping(prs, page_no):
    slide = new_slide(prs)
    add_brandlogy(slide, page_no=page_no, eyebrow="15 · SSOD ↔ SoX 매핑")

    add_textbox(
        slide, SAFE_X, Inches(1.0), CONTENT_W, Inches(0.7),
        "Sync with Self · Other · Data",
        font=FONT_MONO, size=22, color=C_INK, bold=True, align=PP_ALIGN.LEFT, letter_spacing=300,
    )

    # 3-row mapping
    rows = [
        ("Self",  "나 자신과 맞추는 일",     "판단 — SoI",       "결정이 내려지는 자리"),
        ("Other", "사람·조직과 맞추는 일",   "일하는 자리 — SoE", "메시지·결재가 흐르는 곳"),
        ("Data",  "데이터와 맞추는 일",      "기록 — SoR",       "회사에 적힌 사실"),
    ]
    row_y = Inches(2.0)
    row_h = Inches(1.3)
    row_gap = Inches(0.18)

    for i, (kw_en, kw_kr, sox, body) in enumerate(rows):
        y = row_y + (row_h + row_gap) * i
        is_last = i == 2
        add_rect(slide, SAFE_X, y, CONTENT_W, row_h, fill=C_SURFACE_CARD if not is_last else C_CANVAS, radius=Inches(0.15),
                 line_color=C_CORAL if is_last else None, line_width=2 if is_last else 0)
        # left mono
        add_textbox(
            slide, SAFE_X + Inches(0.5), y + Inches(0.3), Inches(2), Inches(0.4),
            kw_en,
            font=FONT_MONO, size=22, color=C_CORAL, bold=True, letter_spacing=200,
        )
        add_textbox(
            slide, SAFE_X + Inches(0.5), y + Inches(0.75), Inches(2), Inches(0.4),
            kw_kr,
            font=FONT_KR, size=14, color=C_MUTED,
        )

        # middle arrow
        add_textbox(
            slide, SAFE_X + Inches(3.5), y + Inches(0.45), Inches(1), Inches(0.5),
            "→",
            font=FONT_MONO, size=24, color=C_MUTED_SOFT, align=PP_ALIGN.CENTER,
        )

        # right
        add_textbox(
            slide, SAFE_X + Inches(5), y + Inches(0.25), CONTENT_W - Inches(5.5), Inches(0.5),
            sox,
            font=FONT_KR, size=22, color=C_INK, bold=True, letter_spacing=-100,
        )
        add_textbox(
            slide, SAFE_X + Inches(5), y + Inches(0.8), CONTENT_W - Inches(5.5), Inches(0.4),
            body,
            font=FONT_KR, size=13, color=C_BODY,
        )

    # footer thesis
    add_textbox(
        slide, SAFE_X, Inches(6.5), CONTENT_W, Inches(0.5),
        "셋을 한 자리에 맞추는 일 — SSOD가 합니다.",
        font=FONT_KR, size=18, color=C_INK, bold=True, align=PP_ALIGN.CENTER, letter_spacing=-50,
    )
    return slide


# --- 16. photo wall (s16)
def slide_photo_wall(prs, page_no):
    slide = new_slide(prs)
    add_brandlogy(slide, page_no=page_no, eyebrow="16 · 함께한 일들")

    add_textbox(
        slide, SAFE_X, Inches(0.95), CONTENT_W, Inches(0.6),
        "이 일을 — 어디서 어떻게 하고 있는지.",
        font=FONT_KR, size=28, color=C_INK, letter_spacing=-150,
    )

    # 4 photo slots (placeholder cards)
    photos = [
        ("Notion Camp HR", "노션 본사 워크숍 · 2025"),
        ("상장사·중견기업", "6개월 상주 컨설팅 · NDA"),
        ("n8n·Notion 공식 행사", "강연 · Seoul · Online"),
        ("자유 슬롯", "발표 현장 photo drop"),
    ]
    grid_y = Inches(2.0)
    cols = 4
    gap = Inches(0.2)
    cell_w = (CONTENT_W - gap * (cols - 1)) / cols
    cell_h = Inches(3.6)  # 4:5 비율 근접

    for i, (title, sub) in enumerate(photos):
        x = SAFE_X + (cell_w + gap) * i
        is_placeholder = i == 3
        add_rect(
            slide, x, grid_y, cell_w, cell_h,
            fill=C_SURFACE_CARD if not is_placeholder else C_CANVAS,
            line_color=C_HAIRLINE if is_placeholder else None,
            line_width=2 if is_placeholder else 0,
            radius=Inches(0.15),
        )
        # placeholder mono icon
        add_textbox(
            slide, x, grid_y + Inches(1.3), cell_w, Inches(0.5),
            "[ photo ]",
            font=FONT_MONO, size=14, color=C_MUTED_SOFT, align=PP_ALIGN.CENTER, letter_spacing=200,
        )
        # caption
        add_textbox(
            slide, x + Inches(0.2), grid_y + cell_h + Inches(0.15), cell_w - Inches(0.4), Inches(0.4),
            title,
            font=FONT_KR, size=14, color=C_INK, bold=True, letter_spacing=-50,
        )
        add_textbox(
            slide, x + Inches(0.2), grid_y + cell_h + Inches(0.55), cell_w - Inches(0.4), Inches(0.3),
            sub,
            font=FONT_MONO, size=10, color=C_MUTED_SOFT, letter_spacing=150,
        )

    # footer thesis
    add_textbox(
        slide, SAFE_X, Inches(7.05), CONTENT_W, Inches(0.4),
        "직군 · 회사 크기 달라도 — 진단 렌즈는 같습니다.",
        font=FONT_KR, size=14, color=C_BODY, align=PP_ALIGN.CENTER,
    )
    return slide


# --- 17. slogan dark (s17)
def slide_slogan_dark(prs, page_no):
    slide = new_slide(prs, dark=True)
    add_brandlogy(slide, page_no=page_no, eyebrow="17 · TAKE-AWAY", dark=True)

    add_textbox(
        slide, SAFE_X, Inches(1.5), CONTENT_W, Inches(0.4),
        "오늘의 한 줄",
        font=FONT_MONO, size=13, color=C_CORAL, bold=True, letter_spacing=350,
    )

    # 큰 slogan
    add_multiline_textbox(
        slide,
        SAFE_X, Inches(2.4),
        CONTENT_W, Inches(2.5),
        [
            {"text": "구조가 없으면,", "font": FONT_KR, "size": 72, "color": C_ON_DARK, "letter_spacing": -350, "spacing": 1.1},
            {"text": "AI도 없습니다.", "font": FONT_KR, "size": 72, "color": C_CORAL, "letter_spacing": -350, "spacing": 1.1, "space_before": 8},
        ]
    )

    # 3-항목 anchor
    add_multiline_textbox(
        slide,
        SAFE_X, Inches(5.6),
        CONTENT_W, Inches(1.6),
        [
            {"text": "데이터는 사장님 노트북에.", "font": FONT_KR, "size": 18, "color": C_ON_DARK_SOFT, "letter_spacing": -50},
            {"text": "일은 카톡방에.", "font": FONT_KR, "size": 18, "color": C_ON_DARK_SOFT, "letter_spacing": -50, "space_before": 4},
            {"text": "판단은 사장님 머릿속에.", "font": FONT_KR, "size": 18, "color": C_ON_DARK_SOFT, "letter_spacing": -50, "space_before": 4},
        ]
    )
    return slide


# --- 18. before/after (s18)
def slide_before_after(prs, page_no):
    slide = new_slide(prs)
    add_brandlogy(slide, page_no=page_no, eyebrow="18 · 사례 — BEFORE → AFTER")

    add_textbox(
        slide, SAFE_X, Inches(0.95), CONTENT_W, Inches(0.5),
        "두 상장사급 회사 — 진단 패턴이 같은 두 곳을 섞었습니다.",
        font=FONT_KR, size=14, color=C_MUTED, letter_spacing=-50,
    )

    # 2-col
    col_w = (CONTENT_W - Inches(0.5)) / 2
    col_h = Inches(5.0)
    col_y = Inches(1.7)

    # BEFORE
    add_rect(slide, SAFE_X, col_y, col_w, col_h, fill=C_SURFACE_CARD, radius=Inches(0.2))
    add_textbox(
        slide, SAFE_X + Inches(0.5), col_y + Inches(0.4), col_w - Inches(1), Inches(0.4),
        "BEFORE",
        font=FONT_MONO, size=13, color=C_MUTED, bold=True, letter_spacing=300,
    )
    add_textbox(
        slide, SAFE_X + Inches(0.5), col_y + Inches(0.9), col_w - Inches(1), Inches(0.8),
        "데이터 뽑을 줄 아는 사람 한 명.",
        font=FONT_KR, size=26, color=C_INK, letter_spacing=-150, spacing=1.2,
    )
    befores = [
        "✓  ERP에 매출·재고·운영 지표 매일 정확히",
        "✓  결재 시스템 안에서 정상 작동",
        "✗  임원 \"어제 어디가 잘 됐지\" 한 줄 → 몇 시간 대기",
        "✗  담당자 매일 출근 후 2-3시간 엑셀 복붙",
        "✗  담당자 휴가 = 임원도 어제를 모름",
    ]
    for i, b in enumerate(befores):
        is_check = b.startswith("✓")
        c = C_BODY if is_check else C_CORAL_ACTIVE
        add_textbox(
            slide, SAFE_X + Inches(0.5), col_y + Inches(2.0) + Inches(0.45) * i, col_w - Inches(1), Inches(0.4),
            b,
            font=FONT_KR, size=14, color=c, spacing=1.2,
        )

    # AFTER
    after_x = SAFE_X + col_w + Inches(0.5)
    add_rect(slide, after_x, col_y, col_w, col_h, fill=C_INK, radius=Inches(0.2))
    add_textbox(
        slide, after_x + Inches(0.5), col_y + Inches(0.4), col_w - Inches(1), Inches(0.4),
        "AFTER — 6개월 상주 후",
        font=FONT_MONO, size=13, color=C_CORAL, bold=True, letter_spacing=300,
    )
    add_textbox(
        slide, after_x + Inches(0.5), col_y + Inches(0.9), col_w - Inches(1), Inches(0.8),
        "기다림이 사라졌어요.",
        font=FONT_KR, size=26, color=C_ON_DARK, letter_spacing=-150, spacing=1.2,
    )
    afters = [
        "→  ERP 숫자 매일 자동 노션 DB로",
        "→  본사·현장·임원 동선 노션 안으로",
        "→  매일 출근 전 어제 지표 임원 화면 자동",
        "→  이상 지점 → 담당 부서에 업무 자동 분배",
        "→  회의 들어가기 전 결정 한 단계 끝",
    ]
    for i, a in enumerate(afters):
        add_textbox(
            slide, after_x + Inches(0.5), col_y + Inches(2.0) + Inches(0.45) * i, col_w - Inches(1), Inches(0.4),
            a,
            font=FONT_KR, size=14, color=C_ON_DARK_SOFT, spacing=1.2,
        )

    # 하단 thesis
    add_textbox(
        slide, SAFE_X, Inches(6.85), CONTENT_W, Inches(0.4),
        "규모는 달라도 진단은 같아요 — 구조가 있는가, 없는가.",
        font=FONT_KR, size=14, color=C_BODY, align=PP_ALIGN.CENTER, letter_spacing=-50,
    )
    return slide


# --- 19. AI 본질 (s19)
def slide_ai_essence(prs, page_no):
    slide = new_slide(prs)
    add_brandlogy(slide, page_no=page_no, eyebrow="19 · AI 본질")

    add_textbox(
        slide, SAFE_X, Inches(1.0), CONTENT_W, Inches(0.4),
        "AI 시대에 진짜 비싼 일",
        font=FONT_MONO, size=13, color=C_CORAL, bold=True, letter_spacing=300,
    )

    add_multiline_textbox(
        slide,
        SAFE_X, Inches(1.8),
        CONTENT_W, Inches(2.6),
        [
            {"text": "AI는 일하는 구조를", "font": FONT_KR, "size": 56, "color": C_INK, "letter_spacing": -300, "spacing": 1.15},
            {"text": "스스로 못 만듭니다.", "font": FONT_KR, "size": 56, "color": C_CORAL, "letter_spacing": -300, "spacing": 1.15, "space_before": 6},
        ]
    )

    # 본문
    add_rect(slide, SAFE_X, Inches(5.0), CONTENT_W, Inches(1.8), fill=C_SURFACE_CARD, radius=Inches(0.2))
    add_multiline_textbox(
        slide,
        SAFE_X + Inches(0.6), Inches(5.25),
        CONTENT_W - Inches(1.2), Inches(1.4),
        [
            {"text": "사람 머릿속 판단 기준 →", "font": FONT_KR, "size": 18, "color": C_MUTED, "letter_spacing": -50},
            {"text": "사람이 사람한테 묻고 정리해서 →", "font": FONT_KR, "size": 18, "color": C_MUTED, "letter_spacing": -50, "space_before": 6},
            {"text": "데이터 옆에 같이 적어 둬야 — AI가 일을 시작합니다.", "font": FONT_KR, "size": 19, "color": C_INK, "bold": True, "letter_spacing": -50, "space_before": 8},
        ]
    )
    return slide


# --- 20. Palantir 모델 (s20) — v9.1 F2 시총 anchor 제거 + 톤 중립화
def slide_palantir_fde(prs, page_no):
    slide = new_slide(prs)
    add_brandlogy(slide, page_no=page_no, eyebrow="20 · 글로벌 reference — 팔란티어")

    # 좌측 큰 punch
    add_textbox(
        slide, SAFE_X, Inches(1.2), Inches(6), Inches(0.4),
        "PALANTIR",
        font=FONT_MONO, size=18, color=C_MUTED, bold=True, letter_spacing=400,
    )
    add_multiline_textbox(
        slide,
        SAFE_X, Inches(1.9),
        Inches(6), Inches(3.5),
        [
            {"text": "엔지니어를", "font": FONT_KR, "size": 44, "color": C_INK, "letter_spacing": -200, "spacing": 1.15},
            {"text": "고객 회사 안에", "font": FONT_KR, "size": 44, "color": C_INK, "letter_spacing": -200, "spacing": 1.15, "space_before": 0},
            {"text": "직접 들여보내는 방식.", "font": FONT_KR, "size": 44, "color": C_CORAL, "letter_spacing": -200, "spacing": 1.15, "space_before": 4},
        ]
    )

    # 우측 카드 — Forward Deployed 모델 (FDE 약어는 deck visual anchor만, 발화 X)
    card_x = SAFE_X + Inches(6.5)
    card_w = CONTENT_W - Inches(6.5)
    add_rect(slide, card_x, Inches(1.2), card_w, Inches(5.4), fill=C_INK, radius=Inches(0.2))

    # 영문 약어는 visual anchor (작게, mono) — 발화 부담 회피
    add_textbox(
        slide, card_x + Inches(0.5), Inches(1.6), card_w - Inches(1), Inches(0.4),
        "Forward Deployed",
        font=FONT_MONO, size=18, color=C_ON_DARK_SOFT, bold=True, letter_spacing=200,
    )
    add_textbox(
        slide, card_x + Inches(0.5), Inches(2.1), card_w - Inches(1), Inches(0.5),
        "회사 안 상주",
        font=FONT_KR, size=24, color=C_ON_DARK, bold=True, letter_spacing=-100,
    )
    add_line(slide, card_x + Inches(0.5), Inches(2.95), card_x + card_w - Inches(0.5), Inches(2.95), color=C_SURFACE_DARK_ELEVATED, width_pt=1)

    add_multiline_textbox(
        slide,
        card_x + Inches(0.5), Inches(3.2),
        card_w - Inches(1), Inches(3.2),
        [
            {"text": "본사 사무실에 앉히지 않고", "font": FONT_KR, "size": 16, "color": C_ON_DARK, "spacing": 1.4},
            {"text": "고객 회사 안에 데려다 놓습니다.", "font": FONT_KR, "size": 16, "color": C_ON_DARK, "spacing": 1.4, "space_before": 6},
            {"text": "어떻게 판단하는지 — ", "font": FONT_KR, "size": 16, "color": C_ON_DARK_SOFT, "spacing": 1.4, "space_before": 12},
            {"text": "옆에서 같이 일하면서", "font": FONT_KR, "size": 16, "color": C_ON_DARK, "spacing": 1.4, "space_before": 4},
            {"text": "데이터 옆에 정리해 둡니다.", "font": FONT_KR, "size": 16, "color": C_ON_DARK, "spacing": 1.4, "space_before": 4},
            {"text": "그게 — 그 회사 방식이고 본질입니다.", "font": FONT_KR, "size": 14, "color": C_CORAL, "bold": False, "spacing": 1.4, "space_before": 14, "letter_spacing": -50},
        ]
    )
    return slide


# --- 21. 같은 방식 — 한국에선 쏘드가 (s21) — v9.1 F3·F8 자기 등급 매김 제거 + 차별화 강화
def slide_ssod_fde(prs, page_no):
    slide = new_slide(prs)
    add_brandlogy(slide, page_no=page_no, eyebrow="21 · 같은 방식 — 한국에선 쏘드가")

    # 큰 헤드라인
    add_textbox(
        slide, SAFE_X, Inches(1.2), CONTENT_W, Inches(0.5),
        "같은 방식 — 회사를 옆에서 직접 짓는 일.",
        font=FONT_KR, size=14, color=C_MUTED, letter_spacing=-50,
    )
    add_multiline_textbox(
        slide,
        SAFE_X, Inches(1.9),
        CONTENT_W, Inches(2.0),
        [
            {"text": "한국에선 — 쏘드가 합니다.", "font": FONT_KR, "size": 56, "color": C_INK, "letter_spacing": -300, "spacing": 1.1},
        ]
    )

    # 비교 표 (Palantir vs SSOD) — v9.1 F8: "핵심" 차별화
    table_y = Inches(3.5)
    table_h = Inches(3.2)
    add_rect(slide, SAFE_X, table_y, CONTENT_W, table_h, fill=C_SURFACE_CARD, radius=Inches(0.2))

    col_w = CONTENT_W / 3
    headers = ["축",         "팔란티어 방식",            "쏘드의 방식"]
    rows = [
        ("자리",       "본사 사무실에 앉히지 않음",  "본사 사무실에 앉히지 않음"),
        ("어디로",     "고객 회사 안에",            "고객 회사 안에"),
        ("기간",       "프로젝트 단위 상주",         "6개월 ~ 1년 상주"),
        ("정리",       "판단 기준 → 데이터 옆에",   "한국 회사 OS 맥락에서 —\n카톡·ERP·노션 결합으로"),
    ]

    # header row
    hy = table_y + Inches(0.3)
    for i, h in enumerate(headers):
        x = SAFE_X + col_w * i
        add_textbox(
            slide, x + Inches(0.3), hy, col_w - Inches(0.3), Inches(0.4),
            h,
            font=FONT_MONO, size=11, color=C_MUTED_SOFT, bold=True, letter_spacing=300,
        )
    # divider
    add_line(slide, SAFE_X + Inches(0.3), hy + Inches(0.55), SAFE_X + CONTENT_W - Inches(0.3), hy + Inches(0.55), color=C_HAIRLINE, width_pt=1.5)

    for r, (axis, p, s) in enumerate(rows):
        ry = hy + Inches(0.7) + Inches(0.55) * r
        # axis
        add_textbox(
            slide, SAFE_X + Inches(0.3), ry, col_w - Inches(0.3), Inches(0.4),
            axis,
            font=FONT_KR, size=14, color=C_MUTED,
        )
        # palantir
        add_textbox(
            slide, SAFE_X + col_w + Inches(0.3), ry, col_w - Inches(0.3), Inches(0.4),
            p,
            font=FONT_KR, size=13, color=C_BODY, spacing=1.3,
        )
        # ssod — coral은 차별화 row 3(정리)만, voltage 보호
        is_emphasis = r == 3
        ssod_color = C_CORAL if is_emphasis else C_INK
        add_textbox(
            slide, SAFE_X + col_w * 2 + Inches(0.3), ry, col_w - Inches(0.3), Inches(0.5),
            s,
            font=FONT_KR, size=13,
            color=ssod_color,
            bold=is_emphasis,
            spacing=1.3,
        )
    return slide


# --- 22. bridge — 한국에선 쏘드가 (s22)
def slide_bridge(prs, page_no):
    slide = new_slide(prs)
    add_brandlogy(slide, page_no=page_no, eyebrow="22 · 도구")

    add_textbox(
        slide, SAFE_X, Inches(2.5), CONTENT_W, Inches(0.5),
        "그래서 — 어떻게 만드냐.",
        font=FONT_KR, size=22, color=C_MUTED, align=PP_ALIGN.CENTER, letter_spacing=-50,
    )

    add_multiline_textbox(
        slide,
        SAFE_X, Inches(3.3),
        CONTENT_W, Inches(2.5),
        [
            {"text": "도구 — 네 가지로.", "font": FONT_KR, "size": 64, "color": C_INK, "letter_spacing": -300, "spacing": 1.1, "align": PP_ALIGN.CENTER},
        ]
    )
    add_textbox(
        slide, SAFE_X, Inches(5.5), CONTENT_W, Inches(0.5),
        "Notion · n8n · Claude · Microsoft Excel",
        font=FONT_MONO, size=16, color=C_CORAL, bold=True, align=PP_ALIGN.CENTER, letter_spacing=300,
    )
    return slide


# --- 23. tool grid 3종 (s23)
def slide_tool_grid(prs, page_no):
    slide = new_slide(prs)
    add_brandlogy(slide, page_no=page_no, eyebrow="23 · 도구 3종 + 1")

    add_textbox(
        slide, SAFE_X, Inches(0.95), CONTENT_W, Inches(0.6),
        "회사의 일하는 구조 — 어떻게 만드는가.",
        font=FONT_KR, size=28, color=C_INK, letter_spacing=-150,
    )

    # 3-card grid
    tools = [
        ("Notion",          "SoE — 일하는 자리",
         "메시지·결재·문서가 한 자리에 모이게."),
        ("Excel → Notion",   "SoR — 데이터 인입",
         "회사 곳곳에 흩어진 엑셀들을 Notion으로."),
        ("Claude",          "SoI — 판단 시스템",
         "사람 머릿속 기준을 코드로 옮겨 적고 자동 판단."),
    ]
    grid_y = Inches(2.3)
    cols = 3
    gap = Inches(0.3)
    cell_w = (CONTENT_W - gap * (cols - 1)) / cols
    cell_h = Inches(3.5)

    for i, (name, frame, body) in enumerate(tools):
        x = SAFE_X + (cell_w + gap) * i
        add_rect(slide, x, grid_y, cell_w, cell_h, fill=C_SURFACE_CARD, radius=Inches(0.2))
        # frame mono on top
        add_textbox(
            slide, x + Inches(0.4), grid_y + Inches(0.4), cell_w - Inches(0.8), Inches(0.3),
            frame,
            font=FONT_MONO, size=11, color=C_CORAL, bold=True, letter_spacing=250,
        )
        # name (large)
        add_textbox(
            slide, x + Inches(0.4), grid_y + Inches(1.0), cell_w - Inches(0.8), Inches(0.8),
            name,
            font=FONT_KR, size=30, color=C_INK, bold=True, letter_spacing=-150,
        )
        # divider
        add_line(slide, x + Inches(0.4), grid_y + Inches(2.0), x + cell_w - Inches(0.4), grid_y + Inches(2.0), color=C_HAIRLINE, width_pt=1)
        # body
        add_textbox(
            slide, x + Inches(0.4), grid_y + Inches(2.2), cell_w - Inches(0.8), Inches(1.2),
            body,
            font=FONT_KR, size=14, color=C_BODY, spacing=1.4,
        )

    # 하단 flow-connector cue
    add_textbox(
        slide, SAFE_X, Inches(6.1), CONTENT_W, Inches(0.4),
        "↓  ↓  ↓",
        font=FONT_MONO, size=22, color=C_CORAL, align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide, SAFE_X, Inches(6.55), CONTENT_W, Inches(0.4),
        "셋을 잇는 자리가 — 다음 슬라이드.",
        font=FONT_KR, size=14, color=C_MUTED, align=PP_ALIGN.CENTER,
    )
    return slide


# --- 24. n8n 혈관 (s24)
def slide_n8n_vessel(prs, page_no):
    slide = new_slide(prs)
    add_brandlogy(slide, page_no=page_no, eyebrow="24 · 혈관")

    # 좌측
    add_textbox(
        slide, SAFE_X, Inches(1.2), Inches(6), Inches(0.4),
        "셋을 잇는 자리",
        font=FONT_MONO, size=14, color=C_MUTED, letter_spacing=300,
    )
    add_textbox(
        slide, SAFE_X, Inches(1.9), Inches(6), Inches(3),
        "n8n",
        font=FONT_KR, size=140, color=C_CORAL, bold=False, letter_spacing=-400,
    )
    add_textbox(
        slide, SAFE_X, Inches(4.6), Inches(6), Inches(0.4),
        "automation vessel",
        font=FONT_MONO, size=14, color=C_MUTED, letter_spacing=300,
    )

    # 우측 vessel diagram
    diag_x = SAFE_X + Inches(6.5)
    diag_w = CONTENT_W - Inches(6.5)

    add_rect(slide, diag_x, Inches(1.2), diag_w, Inches(5.4), fill=C_INK, radius=Inches(0.2))

    # 3 axes inside dark
    axes = [("SoR", "데이터"), ("SoE", "일하는 자리"), ("SoI", "판단")]
    ay = Inches(1.7)
    ah = Inches(0.9)
    agap = Inches(0.2)
    for i, (f, role) in enumerate(axes):
        y = ay + (ah + agap) * i
        add_rect(slide, diag_x + Inches(0.5), y, diag_w - Inches(1), ah, fill=C_SURFACE_DARK_ELEVATED, radius=Inches(0.12))
        add_textbox(
            slide, diag_x + Inches(0.8), y + Inches(0.2), Inches(1.5), Inches(0.4),
            f,
            font=FONT_MONO, size=18, color=C_CORAL, bold=True, letter_spacing=300,
        )
        add_textbox(
            slide, diag_x + Inches(2.4), y + Inches(0.25), diag_w - Inches(2.9), Inches(0.4),
            role,
            font=FONT_KR, size=16, color=C_ON_DARK, bold=True, letter_spacing=-50,
        )

    # connector vessel bar
    vessel_y = Inches(5.2)
    add_rect(slide, diag_x + Inches(0.5), vessel_y, diag_w - Inches(1), Inches(0.5), fill=C_CORAL, radius=Inches(0.08))
    add_textbox(
        slide, diag_x + Inches(0.5), vessel_y + Inches(0.1), diag_w - Inches(1), Inches(0.3),
        "n8n  ·  data flows through",
        font=FONT_MONO, size=12, color=C_ON_PRIMARY, bold=True, align=PP_ALIGN.CENTER, letter_spacing=250,
    )

    # 하단 thesis
    add_textbox(
        slide, diag_x + Inches(0.5), Inches(5.9), diag_w - Inches(1), Inches(0.5),
        "이 네 가지로 — 회사의 일하는 구조를 만듭니다.",
        font=FONT_KR, size=13, color=C_ON_DARK_SOFT, align=PP_ALIGN.CENTER,
    )
    return slide


# --- 25. closing thesis (s25)
def slide_closing_thesis(prs, page_no):
    slide = new_slide(prs, dark=True)
    add_brandlogy(slide, page_no=page_no, eyebrow="25 · CLOSING THESIS", dark=True)

    # 큰 슬로건
    add_multiline_textbox(
        slide,
        SAFE_X, Inches(1.8),
        CONTENT_W, Inches(3.0),
        [
            {"text": "구조가 없으면,", "font": FONT_KR, "size": 80, "color": C_ON_DARK, "letter_spacing": -400, "spacing": 1.1, "align": PP_ALIGN.CENTER},
            {"text": "AI도 없습니다.", "font": FONT_KR, "size": 80, "color": C_CORAL, "letter_spacing": -400, "spacing": 1.1, "align": PP_ALIGN.CENTER, "space_before": 8},
        ]
    )

    # 부속
    add_textbox(
        slide, SAFE_X, Inches(5.6), CONTENT_W, Inches(0.6),
        "그 구조를 — SSOD가 만듭니다.",
        font=FONT_KR, size=28, color=C_ON_DARK, align=PP_ALIGN.CENTER, letter_spacing=-150,
    )
    add_textbox(
        slide, SAFE_X, Inches(6.3), CONTENT_W, Inches(0.5),
        "사람이 사람 옆에서 같이 — 데이터·일·판단을 한 자리에.",
        font=FONT_KR, size=16, color=C_ON_DARK_SOFT, align=PP_ALIGN.CENTER, letter_spacing=-50,
    )
    return slide


# --- 26. 어느 회사든 — 두 길 (s26, v9 s26·s27 통합 NEW)
# v9.1 F1: 100인+ "전담" 자랑 톤 제거, "어느 회사든" inclusion 톤. duál model → 라우팅 framing.
# v9.1 F9: Academy 잠재 고객 자존감 보호 카피 ("자격 회사만" → "복잡해지는 곳").
def slide_dual_path(prs, page_no):
    slide = new_slide(prs)
    add_brandlogy(slide, page_no=page_no, eyebrow="26 · 앞으로 — 두 길로")

    # eyebrow + headline (inclusion 톤)
    add_textbox(
        slide, SAFE_X, Inches(1.0), CONTENT_W, Inches(0.4),
        "어느 회사든 — 같은 구조를 깔 수 있게",
        font=FONT_KR, size=14, color=C_MUTED, letter_spacing=-50,
    )
    add_multiline_textbox(
        slide,
        SAFE_X, Inches(1.6),
        CONTENT_W, Inches(1.8),
        [
            {"text": "쏘드는 — 두 길로 갑니다.", "font": FONT_KR, "size": 52, "color": C_INK, "letter_spacing": -300, "spacing": 1.1},
        ]
    )

    # 2-tier — 큰 회사 / 작은 회사 (3-tier 압축 → 2-tier inclusion)
    tier_y = Inches(3.6)
    tier_h = Inches(2.9)
    cols = 2
    gap = Inches(0.4)
    cell_w = (CONTENT_W - gap) / 2

    # Tier 1 — 큰 회사 (dark, 본업)
    x1 = SAFE_X
    add_rect(slide, x1, tier_y, cell_w, tier_h, fill=C_INK, radius=Inches(0.2))
    add_textbox(
        slide, x1 + Inches(0.5), tier_y + Inches(0.4), cell_w - Inches(1), Inches(0.3),
        "큰 회사",
        font=FONT_MONO, size=11, color=C_CORAL, bold=True, letter_spacing=300,
    )
    add_textbox(
        slide, x1 + Inches(0.5), tier_y + Inches(0.85), cell_w - Inches(1), Inches(0.7),
        "회사 안에 들어가서, 같이 일하면서.",
        font=FONT_KR, size=22, color=C_ON_DARK, bold=True, letter_spacing=-150, spacing=1.2,
    )
    add_textbox(
        slide, x1 + Inches(0.5), tier_y + Inches(1.85), cell_w - Inches(1), Inches(0.3),
        "100명 ~ · 디지털 전환 + AI 전환이 복잡해지는 곳",
        font=FONT_KR, size=13, color=C_ON_DARK_SOFT, letter_spacing=-30,
    )
    add_line(slide, x1 + Inches(0.5), tier_y + Inches(2.2), x1 + cell_w - Inches(0.5), tier_y + Inches(2.2), color=C_SURFACE_DARK_ELEVATED, width_pt=1)
    add_textbox(
        slide, x1 + Inches(0.5), tier_y + Inches(2.35), cell_w - Inches(1), Inches(0.5),
        "쏘드의 본업 · 상주 컨설팅",
        font=FONT_KR, size=14, color=C_CORAL, bold=True, letter_spacing=-50,
    )

    # Tier 2 — 작은 회사 (cream, 듀얼)
    x2 = SAFE_X + cell_w + gap
    add_rect(slide, x2, tier_y, cell_w, tier_h, fill=C_SURFACE_CARD, radius=Inches(0.2))
    add_textbox(
        slide, x2 + Inches(0.5), tier_y + Inches(0.4), cell_w - Inches(1), Inches(0.3),
        "더 작은 회사",
        font=FONT_MONO, size=11, color=C_MUTED, bold=True, letter_spacing=300,
    )
    add_textbox(
        slide, x2 + Inches(0.5), tier_y + Inches(0.85), cell_w - Inches(1), Inches(0.7),
        "쏘드가 만든 구조를 가져가서.",
        font=FONT_KR, size=22, color=C_INK, bold=True, letter_spacing=-150, spacing=1.2,
    )
    add_textbox(
        slide, x2 + Inches(0.5), tier_y + Inches(1.85), cell_w - Inches(1), Inches(0.3),
        "Notion 시스템 패키지  +  직접 만드는 법 강의",
        font=FONT_KR, size=13, color=C_BODY, letter_spacing=-30,
    )
    add_line(slide, x2 + Inches(0.5), tier_y + Inches(2.2), x2 + cell_w - Inches(0.5), tier_y + Inches(2.2), color=C_HAIRLINE, width_pt=1)
    # vapor → 현재진행형
    add_textbox(
        slide, x2 + Inches(0.5), tier_y + Inches(2.35), cell_w - Inches(1), Inches(0.5),
        "지금 만들고 있습니다 · 자기 회사에 직접 적용",
        font=FONT_KR, size=13, color=C_MUTED, letter_spacing=-30,
    )

    # 하단 footer
    add_textbox(
        slide, SAFE_X, Inches(6.75), CONTENT_W, Inches(0.4),
        "상주는 쏘드가, 패키지·강의는 — 누구나. 한국 회사의 일하는 구조를 같이 짓겠습니다.",
        font=FONT_KR, size=14, color=C_INK, align=PP_ALIGN.CENTER, letter_spacing=-50,
    )
    return slide


# --- 27. referral + 인사 (v9 s28 → v9.1 s27) — v9.1 F5 hook 단일화
def slide_referral_closing(prs, page_no):
    slide = new_slide(prs, dark=True)
    add_brandlogy(slide, page_no=page_no, eyebrow="27 · REFERRAL", dark=True)

    # eyebrow
    add_textbox(
        slide, SAFE_X, Inches(1.0), CONTENT_W, Inches(0.4),
        "오늘의 부탁 한 줄",
        font=FONT_MONO, size=13, color=C_CORAL, bold=True, letter_spacing=400,
    )

    # v9.1 F5 — single mental hook (1 trigger phrase + routing 분리)
    add_multiline_textbox(
        slide,
        SAFE_X, Inches(1.7),
        CONTENT_W, Inches(2.6),
        [
            {"text": "거래처·지인 중에", "font": FONT_KR, "size": 30, "color": C_ON_DARK_SOFT, "letter_spacing": -100, "spacing": 1.3},
            {"text": "ERP는 있는데 엑셀에 의존하는 회사 —", "font": FONT_KR, "size": 38, "color": C_ON_DARK, "letter_spacing": -150, "spacing": 1.3, "space_before": 6, "bold": True},
            {"text": "저를 떠올려주세요.", "font": FONT_KR, "size": 48, "color": C_CORAL, "letter_spacing": -250, "spacing": 1.3, "space_before": 10},
        ]
    )

    # 부가 단서 (작게)
    add_textbox(
        slide, SAFE_X, Inches(4.95), CONTENT_W, Inches(0.4),
        "데이터·일·판단이 사장님 머릿속에만 있는 회사면 — 더 정확하게 매칭됩니다.",
        font=FONT_KR, size=14, color=C_ON_DARK_SOFT, letter_spacing=-50,
    )

    # 라우팅 mini grid — 2-tier 단순화 (v9.1: 큰/작은 회사)
    route_y = Inches(5.6)
    route_h = Inches(0.85)
    route_cards = [
        ("큰 회사",   "→  상주 컨설팅"),
        ("작은 회사", "→  Notion 시스템 패키지·강의"),
    ]
    cols = 2
    gap = Inches(0.3)
    cell_w = (CONTENT_W - gap) / 2
    for i, (target, route) in enumerate(route_cards):
        x = SAFE_X + (cell_w + gap) * i
        add_rect(slide, x, route_y, cell_w, route_h, fill=C_SURFACE_DARK_ELEVATED, radius=Inches(0.1))
        add_textbox(
            slide, x + Inches(0.4), route_y + Inches(0.12), cell_w - Inches(0.8), Inches(0.3),
            target,
            font=FONT_MONO, size=11, color=C_ON_DARK_SOFT, letter_spacing=200,
        )
        add_textbox(
            slide, x + Inches(0.4), route_y + Inches(0.42), cell_w - Inches(0.8), Inches(0.4),
            route,
            font=FONT_KR, size=15, color=C_ON_DARK, bold=True, letter_spacing=-50,
        )

    # 인사 footer
    add_textbox(
        slide, SAFE_X, Inches(6.75), CONTENT_W, Inches(0.4),
        "회사의 일하는 구조를 만드는 일 — 쏘드 컴퍼니의 김지명이었습니다. 감사합니다.",
        font=FONT_KR, size=13, color=C_ON_DARK_SOFT, align=PP_ALIGN.CENTER, letter_spacing=-50,
    )
    return slide


# ============================================================
# 빌드
# ============================================================

def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # v9.1: 27 slides (v9 28장 → s26·s27 통합으로 1장 감)
    slide_cover(prs, "01")              # 1
    slide_intro_centered(prs, "02")     # 2
    slide_logo_wall(prs, "03")          # 3
    slide_ambassador_punch(prs, "04")   # 4
    slide_activity(prs, "05")           # 5
    slide_opening_hook(prs, "06")       # 6

    # ③ SoR/SoE/SoI 정의·통찰 ×3
    # v9.1 F4: coral 강조는 s11 SoI 정의(thesis 진입) + s12 SoI 통찰만
    slide_so_def(prs, "07", frame="SoR", frame_letter="Record",
                 no="07", name_kr="데이터 — 회사의 진짜 사실",
                 body="어제 누가 뭘 샀고, 재고가 얼마고, 단골이 누군지. 회사에 사실로 적혀 있는 모든 것.",
                 examples="ERP · 회계 시스템 · 고객 명단 · 매출·재고·운영 지표")
    slide_so_insight(prs, "08", frame="SoR",
                     no="08",
                     headline="거래 데이터만으로는 부족합니다.",
                     body="\"우리 회사가 어떻게 판단하는지\"도 같이 적혀 있어야 합니다.\n이게 빠진 회사가 가장 많아요. 세 번째에서 다시 잡겠습니다.")
    slide_so_def(prs, "09", frame="SoE", frame_letter="Engagement",
                 no="09", name_kr="일하는 자리 — 사람이 일하는 공간",
                 body="직원들이 진짜로 일하는 공간. 메시지가 오가고, 결정이 내려지는 곳.",
                 examples="카톡방 · 슬랙 · 노션 · 이메일 · 결재 시스템")
    slide_so_insight(prs, "10", frame="SoE",
                     no="10",
                     headline="AI는 이 공간 안으로 — 들어와야 합니다.",
                     body="AI 에이전트가 일하려면, 사람들이 이미 모여 있는 자리에서 호출되고 응답해야 합니다.\n\"AI 콘솔 따로 열어서 명령\" 모델은 무너집니다.")
    slide_so_def(prs, "11", frame="SoI", frame_letter="Intelligence",
                 no="11", name_kr="판단 — 결정이 내려지는 자리",
                 body="데이터를 보고 결정으로 바꿔주는 자리. 옛날엔 사장님이었고, 지금은 대시보드·AI 에이전트가 그 역할을 합니다.",
                 examples="대시보드 · AI 에이전트 · 사장님·임원의 직관",
                 is_climax_frame=True)
    slide_so_insight(prs, "12", frame="SoI",
                     no="12",
                     headline="사람이 판단할 때 — 룰은 그 사람 머릿속에.",
                     body="AI가 그 자리를 받으려면? 그 룰이 데이터 옆에 같이 적혀 있어야 합니다.\n데이터만 있고 룰이 없으면 AI는 아무 판단이나 내요.",
                     is_thesis=True)

    # ③.3 AI 시대 3축 연결
    slide_3axis_connect(prs, "13")

    # ③.5 SSOD
    slide_ssod_reveal(prs, "14")
    slide_ssod_mapping(prs, "15")

    # ③.7 photo wall
    slide_photo_wall(prs, "16")

    # ④ 사례
    slide_slogan_dark(prs, "17")
    slide_before_after(prs, "18")

    # ⑤ AI 본질 + Palantir 모델
    slide_ai_essence(prs, "19")
    slide_palantir_fde(prs, "20")
    slide_ssod_fde(prs, "21")

    # ⑤.7 도구
    slide_bridge(prs, "22")
    slide_tool_grid(prs, "23")
    slide_n8n_vessel(prs, "24")

    # ⑥ 클로징 + 듀얼 모델 (v9.1: s26+s27 통합 → slide_dual_path)
    slide_closing_thesis(prs, "25")
    slide_dual_path(prs, "26")         # v9.1 통합 NEW
    slide_referral_closing(prs, "27")  # v9 s28 → v9.1 s27

    prs.save(str(OUT_PATH))
    print(f"OK: {OUT_PATH}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
