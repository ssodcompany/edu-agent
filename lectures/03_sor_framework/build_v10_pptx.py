"""BNI v10 PPT builder — Monochrome Punch (14 slides, 검정+흰색 only).

Spec: design-v10-spec.md (2026-05-25 SHIP) + photo wall s4 추가 (2026-05-26).
Run: python build_v10_pptx.py
Output: BNI_SoR_v10.pptx (16:9, 1280x720)
"""
from pathlib import Path

from lxml import etree

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

ROOT = Path(__file__).parent
ASSETS = ROOT / "assets"
PHOTOS = ASSETS / "photos"
OUTPUT = ROOT / "BNI_SoR_v10.pptx"

CANVAS_W = Inches(13.333)
CANVAS_H = Inches(7.5)
TOTAL_SLIDES = 14

# Linear DESIGN.md 이식 (alpha v10.1, monochrome 유지 + ink·spacing refinement)
BG_BLACK = RGBColor(0x01, 0x01, 0x02)  # Linear canvas — near-black with faint blue tint
INK_WHITE = RGBColor(0xF7, 0xF8, 0xF8)  # Linear ink — light gray (덜 가혹)
MUTED_GRAY = RGBColor(0x8A, 0x8F, 0x98)  # Linear ink-subtle (caption)
INK_TERTIARY = RGBColor(0x62, 0x66, 0x6D)  # Linear ink-tertiary (page no, eyebrow)
HAIRLINE = RGBColor(0x23, 0x25, 0x2A)  # 1px border / divider
SURFACE_1 = RGBColor(0x0F, 0x10, 0x11)  # card lift (diagram 박스 fill)

LAVENDER = RGBColor(0x5E, 0x6A, 0xD2)  # Linear signature accent (brand mark only)

FONT_BLACK = "Pretendard Black"
FONT_LIGHT = "Pretendard Light"
FONT_MONO = "SF Mono"

PUNCH_SPACING_PCT = -3.5  # Linear display aggressive negative tracking


# ────────────────────────────────────────────────────────────
# Helpers
# ────────────────────────────────────────────────────────────


def set_font(run, font_name: str):
    """python-pptx font.name은 Latin만 잡음. Hangul/CJK 위해 ea·cs track도 동일 폰트로 강제."""
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = etree.SubElement(rPr, qn(tag))
        el.set("typeface", font_name)


def set_letter_spacing(run, spacing_pct: float):
    """Letter-spacing in percent of font size. Negative = tighter (Linear signature).

    OOXML spc unit = 1/100 pt. 200pt × -3.5% = -7pt = spc='-700'.
    """
    if run.font.size is None:
        return
    pt = run.font.size.pt
    spc_val = int(pt * spacing_pct)
    rPr = run._r.get_or_add_rPr()
    rPr.set("spc", str(spc_val))


def add_black_slide(prs: Presentation):
    """16:9 빈 슬라이드 + 검정 배경."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, CANVAS_W, CANVAS_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_BLACK
    bg.line.fill.background()
    return slide


def add_gradient_slide(prs: Presentation):
    """16:9 빈 슬라이드 + diagonal 그라데이션 (좌상 검정 → 보라 → 우하 라일락).

    3-stop gradient, 135° diagonal (top-left → bottom-right).
    """
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, CANVAS_W, CANVAS_H)
    bg.line.fill.background()

    sp_pr = bg.fill._xPr
    for tag in ("a:solidFill", "a:gradFill", "a:noFill", "a:pattFill", "a:blipFill"):
        el = sp_pr.find(qn(tag))
        if el is not None:
            sp_pr.remove(el)

    grad = etree.SubElement(sp_pr, qn("a:gradFill"))
    grad.set("flip", "none")
    grad.set("rotWithShape", "1")
    gs_lst = etree.SubElement(grad, qn("a:gsLst"))

    # 3-stop: 좌상 검정 영역 넓게 → 중앙 진보라 → 우하 라일락
    # 검정 0~35%, 보라 70%, 라일락 100% — reference 톤 매칭
    stops = [
        ("0", "000000"),
        ("35000", "000000"),
        ("70000", "6a5acd"),
        ("100000", "e8b0e8"),
    ]
    for pos, hex_color in stops:
        gs = etree.SubElement(gs_lst, qn("a:gs"))
        gs.set("pos", pos)
        srgb = etree.SubElement(gs, qn("a:srgbClr"))
        srgb.set("val", hex_color)

    lin = etree.SubElement(grad, qn("a:lin"))
    lin.set("ang", "2700000")  # 45° = 좌상→우하 diagonal
    lin.set("scaled", "0")

    return slide


def add_punch_text(
    slide,
    text: str,
    font_size_pt: int,
    *,
    font_name: str = FONT_BLACK,
    color: RGBColor = INK_WHITE,
    top_inches: float = 2.8,
    height_inches: float = 2.0,
):
    """슬라이드 중앙 가로 정렬 한 줄 텍스트."""
    tb = slide.shapes.add_textbox(
        Inches(0), Inches(top_inches), CANVAS_W, Inches(height_inches)
    )
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    set_font(run, font_name)
    run.font.size = Pt(font_size_pt)
    run.font.color.rgb = color
    # Linear aggressive negative tracking on display punch
    if font_name == FONT_BLACK and font_size_pt >= 80:
        set_letter_spacing(run, PUNCH_SPACING_PCT)
    return tb


def add_caption(
    slide,
    text: str,
    *,
    top_inches: float = 5.6,
    font_size_pt: int = 18,
    color: RGBColor = MUTED_GRAY,
    align=PP_ALIGN.CENTER,
):
    """보조 회색 캡션."""
    tb = slide.shapes.add_textbox(
        Inches(0), Inches(top_inches), CANVAS_W, Inches(0.5)
    )
    p = tb.text_frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    set_font(r, FONT_LIGHT)
    r.font.size = Pt(font_size_pt)
    r.font.color.rgb = color
    return tb


SSOD_LOGO = ASSETS / "ssod-wordmark.png"  # 가로형 워드마크 (마크 + SSOD 영문)
SSOD_LOGO_H = Inches(0.62)
SSOD_LOGO_W = Inches(0.62 * 2440 / 1585)  # 원본 비율 2440x1585

S10_PYRAMID_SVG = ASSETS / "s10_pyramid.svg"
S10_PYRAMID_PNG = ASSETS / "s10_pyramid.png"

S1_BG_SVG = ASSETS / "s1_bg_pattern.svg"
S1_BG_PNG = ASSETS / "s1_bg_pattern.png"


def _ensure_svg_png(svg_path, png_path, width, height, background=None):
    """SVG → PNG 변환 (mtime stale-check)."""
    if not svg_path.exists():
        return
    if png_path.exists() and png_path.stat().st_mtime >= svg_path.stat().st_mtime:
        return
    import cairosvg
    kwargs = dict(
        url=str(svg_path), write_to=str(png_path),
        output_width=width, output_height=height,
    )
    if background:
        kwargs["background_color"] = background
    cairosvg.svg2png(**kwargs)


def _ensure_s10_png():
    _ensure_svg_png(S10_PYRAMID_SVG, S10_PYRAMID_PNG, 1500, 1500, "#010102")


def _ensure_s1_bg_png():
    _ensure_svg_png(S1_BG_SVG, S1_BG_PNG, 2666, 1500)


LOGO_WHITE_DIR = ASSETS / "case-logo-white"


def _ensure_white_logo(src_path):
    """검정 BG 통일 — alpha mask만 사용, RGB는 ink-white로 일괄 변환.

    홈페이지 ClientLogoSection `brightness-0 invert opacity-70` CSS를 PIL로 재현.
    SVG 입력은 cairosvg로 PNG 중간 변환 후 동일 처리.
    캐시는 case-logo-white/. 원본 mtime 기반 stale-check.
    """
    LOGO_WHITE_DIR.mkdir(exist_ok=True)

    # SVG 입력 → BG rect 제거 → PNG 중간 변환 (case-logo-white/_svg2png/ 캐시)
    if src_path.suffix.lower() == ".svg":
        import re
        svg_clean_dir = LOGO_WHITE_DIR / "_svg_clean"
        svg_clean_dir.mkdir(exist_ok=True)
        cleaned_svg = svg_clean_dir / src_path.name
        if not cleaned_svg.exists() or cleaned_svg.stat().st_mtime < src_path.stat().st_mtime:
            content = src_path.read_text(encoding="utf-8")
            vb_match = re.search(r'viewBox="([\d.\-\s]+)"', content)
            if vb_match:
                vb = [float(x) for x in vb_match.group(1).split()]
                if len(vb) == 4:
                    _, _, vb_w, vb_h = vb

                    def _remove_bg(m):
                        s = m.group(0)
                        wm = re.search(r'width="([\d.\-]+)"', s)
                        hm = re.search(r'height="([\d.\-]+)"', s)
                        if wm and hm:
                            if float(wm.group(1)) / vb_w > 0.9 and float(hm.group(1)) / vb_h > 0.9:
                                return ""
                        return s

                    content = re.sub(r'<rect[^/>]*/\s*>', _remove_bg, content)
            cleaned_svg.write_text(content, encoding="utf-8")
        svg_png_dir = LOGO_WHITE_DIR / "_svg2png"
        svg_png_dir.mkdir(exist_ok=True)
        png_temp = svg_png_dir / f"{src_path.stem}.png"
        if not png_temp.exists() or png_temp.stat().st_mtime < cleaned_svg.stat().st_mtime:
            import cairosvg
            cairosvg.svg2png(
                url=str(cleaned_svg), write_to=str(png_temp), output_width=1200
            )
        src_path = png_temp

    dst = LOGO_WHITE_DIR / f"{src_path.stem}.png"
    if dst.exists() and dst.stat().st_mtime >= src_path.stat().st_mtime:
        return dst
    from PIL import Image
    import numpy as np
    img = Image.open(src_path).convert("RGBA")
    arr = np.array(img).astype(np.float32)
    a = arr[:, :, 3]
    rgb = arr[:, :, :3]
    # luminance (Rec. 601) — 어두울수록 1.0, 밝을수록 0.0
    lum = 0.299 * rgb[:, :, 0] + 0.587 * rgb[:, :, 1] + 0.114 * rgb[:, :, 2]
    darkness = 1.0 - (lum / 255.0)
    # brightness-0 invert 재현: 검정→흰색 진하게, 흰색→투명
    # original alpha와 darkness weight 곱해서 최종 alpha 계산
    h, w = a.shape
    out = np.zeros((h, w, 4), dtype=np.uint8)
    out[:, :, 0] = 0xF7
    out[:, :, 1] = 0xF8
    out[:, :, 2] = 0xF8
    out[:, :, 3] = np.clip(a * darkness * 0.92, 0, 255).astype(np.uint8)
    img_out = Image.fromarray(out, "RGBA")
    # bbox crop — opacity 있는 영역만 남기고 padding 제거 (cupora 등 캔버스 내 padding 큰 로고 fit 정확도 향상)
    bbox = img_out.getbbox()
    if bbox:
        img_out = img_out.crop(bbox)
    img_out.save(dst)
    return dst


def _fit_logo(slide, path, cell_left_in, cell_top_in, cell_w_in, cell_h_in,
              max_w_pct=0.78, max_h_pct=0.65):
    """로고를 cell 안 center에 비율 유지 fit (brand wall용).

    PIL로 원본 비율 측정 → max_w·max_h에 fit하는 scale 계산 → center align.
    """
    from PIL import Image
    img = Image.open(path)
    iw, ih = img.size
    ratio = iw / ih
    max_w = cell_w_in * max_w_pct
    max_h = cell_h_in * max_h_pct
    if max_w / ratio <= max_h:
        w = max_w
        h = max_w / ratio
    else:
        h = max_h
        w = max_h * ratio
    left = cell_left_in + (cell_w_in - w) / 2
    top = cell_top_in + (cell_h_in - h) / 2
    slide.shapes.add_picture(
        str(path), Inches(left), Inches(top),
        width=Inches(w), height=Inches(h),
    )

# Footer 레이아웃 토큰
FOOTER_DIVIDER_Y = Inches(6.72)  # divider 위치 (footer 영역 확장)
FOOTER_BASELINE_Y = Inches(6.82)  # chip·로고 top baseline


def add_brandlogy(slide, slide_no: int):
    """좌하 NN/14 chip + 하단 hairline divider + 우하 SSOD 가로형 워드마크.

    Footer 시각 균형: divider 위(콘텐츠) / divider 아래(footer 영역) 분리.
    좌측 chip(작음, 보조 정보) vs 우측 워드마크(큼, 브랜드 시그너처) 위계 명확.
    """
    # 하단 hairline divider — footer 영역 윤곽
    divider = slide.shapes.add_connector(
        1, Inches(0.4), FOOTER_DIVIDER_Y, CANVAS_W - Inches(0.4), FOOTER_DIVIDER_Y
    )
    divider.line.color.rgb = HAIRLINE
    divider.line.width = Pt(0.5)

    # 좌하: pill chip — surface-1 fill + hairline border + rounded pill
    chip_w = Inches(1.05)
    chip_h = Inches(0.32)
    # 로고 vertical center와 align (로고 height 0.62, top 6.82 → center 7.13)
    # chip center 7.13 → chip top 7.13 - 0.16 = 6.97
    chip = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(6.97), chip_w, chip_h
    )
    chip.fill.solid()
    chip.fill.fore_color.rgb = SURFACE_1
    chip.line.color.rgb = HAIRLINE
    chip.line.width = Pt(0.75)
    chip.adjustments[0] = 0.5  # pill
    tf = chip.text_frame
    tf.word_wrap = False  # 한 줄 강제
    tf.margin_left = tf.margin_right = Inches(0.1)
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = f"{slide_no:02d} / {TOTAL_SLIDES}"
    set_font(r, FONT_MONO)
    r.font.size = Pt(10)
    r.font.color.rgb = MUTED_GRAY
    set_letter_spacing(r, 4.0)  # eyebrow positive tracking

    # 우하: SSOD 가로형 워드마크 (footer 영역 안, divider 아래)
    if SSOD_LOGO.exists():
        right_margin = Inches(0.4)
        # 로고 top 6.82, height 0.62 → bottom 7.44 (캔버스 7.5 안, divider 6.72 아래)
        top = FOOTER_BASELINE_Y
        left = CANVAS_W - SSOD_LOGO_W - right_margin
        slide.shapes.add_picture(
            str(SSOD_LOGO), left, top, width=SSOD_LOGO_W, height=SSOD_LOGO_H
        )


# ────────────────────────────────────────────────────────────
# Slide builders
# ────────────────────────────────────────────────────────────


def build_s1_cover(prs):
    """s1 표지 — 메인 페이지 강화판.

    레이어: 4-corner L자 cut mark (measurement 신호, CoreSystemGraphic 노하우 echo)
          + lead-in mono "03 · SOR FRAMEWORK" (정보 위계 1단)
          + eyebrow "AI 시대" + punch "준비해야하는 일의 구조"
          + 라벤더 hairline accent (punch 아래 중앙, Linear signature spot)
          + 중앙 SSOD 워드마크 + 하단 발표 meta.
    """
    slide = add_black_slide(prs)

    # ── BG geometric pattern (grid hairline + 라벤더 dot accent + 중앙 fade)
    _ensure_s1_bg_png()
    if S1_BG_PNG.exists():
        slide.shapes.add_picture(
            str(S1_BG_PNG), 0, 0, width=CANVAS_W, height=CANVAS_H
        )

    # ── 4-corner L자 cut mark (페이지 모서리 등록 마크, 컨설팅 measurement 톤)
    cut_len = Inches(0.28)
    cut_margin = Inches(0.45)
    corners = [
        (cut_margin, cut_margin, 1, 1),                                   # top-left
        (CANVAS_W - cut_margin, cut_margin, -1, 1),                       # top-right
        (cut_margin, CANVAS_H - cut_margin, 1, -1),                       # bottom-left
        (CANVAS_W - cut_margin, CANVAS_H - cut_margin, -1, -1),           # bottom-right
    ]
    for x, y, sx, sy in corners:
        # horizontal arm
        h = slide.shapes.add_connector(1, x, y, x + cut_len * sx, y)
        h.line.color.rgb = INK_TERTIARY
        h.line.width = Pt(0.75)
        # vertical arm
        v = slide.shapes.add_connector(1, x, y, x, y + cut_len * sy)
        v.line.color.rgb = INK_TERTIARY
        v.line.width = Pt(0.75)

    # ── lead-in mono (정보 위계 — 강의 시리즈·번호)
    lead = slide.shapes.add_textbox(Inches(0), Inches(1.9), CANVAS_W, Inches(0.4))
    lp = lead.text_frame.paragraphs[0]
    lp.alignment = PP_ALIGN.CENTER
    lr = lp.add_run()
    lr.text = "03  ·  SOR FRAMEWORK"
    set_font(lr, FONT_MONO)
    lr.font.size = Pt(11)
    lr.font.color.rgb = INK_TERTIARY
    set_letter_spacing(lr, 22.0)

    # ── eyebrow: AI 시대
    eyebrow = slide.shapes.add_textbox(Inches(0), Inches(2.55), CANVAS_W, Inches(0.5))
    ep = eyebrow.text_frame.paragraphs[0]
    ep.alignment = PP_ALIGN.CENTER
    er = ep.add_run()
    er.text = "AI 시대"
    set_font(er, FONT_LIGHT)
    er.font.size = Pt(28)
    er.font.color.rgb = MUTED_GRAY

    # ── punch: 준비해야하는 일의 구조
    add_punch_text(
        slide,
        "준비해야하는 일의 구조",
        72,
        top_inches=3.2,
        height_inches=1.6,
    )

    # ── 라벤더 hairline accent (punch 아래 중앙, Linear signature)
    div_w = Inches(0.9)
    div_y = Inches(5.05)
    div_left = (CANVAS_W - div_w) / 2
    divider = slide.shapes.add_connector(1, div_left, div_y, div_left + div_w, div_y)
    divider.line.color.rgb = LAVENDER
    divider.line.width = Pt(1.5)

    # ── 중앙 SSOD 워드마크 (표지 시그너처)
    if SSOD_LOGO.exists():
        h = Inches(0.7)
        w = Inches(0.7 * 2440 / 1585)
        left = (CANVAS_W - w) / 2
        top = Inches(5.55)
        slide.shapes.add_picture(str(SSOD_LOGO), left, top, width=w, height=h)

    # ── 하단 발표 meta (mono small caps)
    meta = slide.shapes.add_textbox(Inches(0), Inches(6.55), CANVAS_W, Inches(0.4))
    mp = meta.text_frame.paragraphs[0]
    mp.alignment = PP_ALIGN.CENTER
    mr = mp.add_run()
    mr.text = "BNI K-CHAPTER   ·   2026.05.27   ·   김지명"
    set_font(mr, FONT_MONO)
    mr.font.size = Pt(10)
    mr.font.color.rgb = INK_TERTIARY
    set_letter_spacing(mr, 18.0)


def build_s2_promise(prs, slide_no=2):
    """s2 인트로 promise — 오늘 가져가실 2가지."""
    slide = add_black_slide(prs)
    add_brandlogy(slide, slide_no)
    add_punch_text(slide, "오늘 가져가실 2가지", 56, top_inches=1.6, height_inches=0.9)

    tb = slide.shapes.add_textbox(Inches(0), Inches(3.3), CANVAS_W, Inches(3.0))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP

    lines = [
        "①  SSOD는 무엇을 하는가?",
        "②  우리 회사는?",
    ]
    for idx, text in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(28)
        r = p.add_run()
        r.text = text
        set_font(r, FONT_LIGHT)
        r.font.size = Pt(40)
        r.font.color.rgb = INK_WHITE


def build_s3_years(prs, slide_no=3):
    """s3 7년 — 메가 punch."""
    slide = add_black_slide(prs)
    add_brandlogy(slide, slide_no)
    add_punch_text(slide, "7년", 240, top_inches=1.8, height_inches=3.5)
    add_caption(slide, "B2B 컨설팅", top_inches=5.4, font_size_pt=24)


def build_s4_dual_ambassador(prs, slide_no=4):
    """s4 국내 유일 + Notion·n8n 배지 (배지만 원색 유지, 가로 중앙 정렬)."""
    slide = add_black_slide(prs)
    add_brandlogy(slide, slide_no)
    add_punch_text(slide, "국내 유일", 130, top_inches=1.4, height_inches=2.2)

    notion_path = ASSETS / "notion-ambassador-badge.png"
    n8n_path = ASSETS / "n8n-amber-ssod-black.png"

    # 실측: notion 400x123 (aspect 3.25), n8n 391x110 (aspect 3.55)
    badge_h_in = 1.1
    badge_h = Inches(badge_h_in)
    notion_w = Inches(badge_h_in * 400 / 123)
    n8n_w = Inches(badge_h_in * 391 / 110)
    gap = Inches(0.45)

    total_w = notion_w + gap + n8n_w
    left_notion = (CANVAS_W - total_w) / 2
    left_n8n = left_notion + notion_w + gap
    top = Inches(4.55)

    if notion_path.exists():
        slide.shapes.add_picture(
            str(notion_path), left_notion, top, width=notion_w, height=badge_h
        )
    if n8n_path.exists():
        slide.shapes.add_picture(
            str(n8n_path), left_n8n, top, width=n8n_w, height=badge_h
        )

    add_caption(
        slide,
        "이중 앰버서더는 국내 SSOD뿐",
        top_inches=6.1,
        font_size_pt=20,
    )


def build_s5_scale(prs, slide_no=5):
    """s5 3 → 상장사 (혼합 색·크기 한 줄)."""
    slide = add_black_slide(prs)
    add_brandlogy(slide, slide_no)

    tb = slide.shapes.add_textbox(Inches(0), Inches(2.8), CANVAS_W, Inches(2.0))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER

    def _run(text, *, size, color):
        r = p.add_run()
        r.text = text
        set_font(r, FONT_BLACK)
        r.font.size = Pt(size)
        r.font.color.rgb = color

    _run("3", size=160, color=INK_WHITE)
    _run("   →   ", size=120, color=MUTED_GRAY)
    _run("상장사", size=140, color=INK_WHITE)

    add_caption(slide, "직원 3명 사무실부터 상장사까지", top_inches=5.6, font_size_pt=20)


def build_s6_dx_ax(prs, slide_no=6):
    """s6 DX · AX — 본업 정의 (SoR/SoE/SoI 진입 transition)."""
    slide = add_black_slide(prs)
    add_brandlogy(slide, slide_no)
    add_punch_text(slide, "DX · AX", 200, top_inches=2.2, height_inches=2.8)
    add_caption(
        slide,
        "Digital Transformation · AI Transformation",
        top_inches=5.4,
        font_size_pt=22,
    )


def build_brand_wall(prs, slide_no=6):
    """brand wall — 9개 partner 로고 3×3 grid (homepage ClientLogoSection 톤).

    "TRUSTED BY" mono eyebrow + grid (각 cell 안 center fit, 비율 보존).
    NDA hidden 1건(아가방) 제외, 홈페이지 partners 리스트 그대로 차용.
    로고는 normalized PNG (이미 white invert 처리됨).
    """
    slide = add_black_slide(prs)
    add_brandlogy(slide, slide_no)

    # ── eyebrow "TRUSTED BY" (homepage 톤 — mono small caps tracking)
    eb = slide.shapes.add_textbox(Inches(0), Inches(1.05), CANVAS_W, Inches(0.4))
    ep = eb.text_frame.paragraphs[0]
    ep.alignment = PP_ALIGN.CENTER
    er = ep.add_run()
    er.text = "TRUSTED BY"
    set_font(er, FONT_MONO)
    er.font.size = Pt(11)
    er.font.color.rgb = INK_TERTIARY
    set_letter_spacing(er, 22.0)

    # ── 16 partners — homepage 9개 + 추가 7개 (분야 다양성: HR·보험·F&B·건축·디자인·헬스·패션)
    partners = [
        # row 1: homepage 신뢰 위계 4개
        "standard-energy.png",   # 스탠다드에너지
        "photoism.png",           # 포토이즘
        "anjeong.png",            # 노무법인 안정
        "모멘텀HR.svg",            # 모멘텀HR — case interview 대기
        # row 2: homepage 4개
        "sonicsleep.png",         # 소닉슬립
        "romer.png",              # 로메르
        "ssalssalssal.png",       # 쌀쌀쌀
        "지산손해사정법인.svg",       # 지산손해사정법인
        # row 3: homepage 3개 + 추가 1개
        "deeplogic.png",          # 딥로직
        "gamdong.png",            # 감동한의원
        "ddoksori.png",           # 똑소리나는부동산
        "쿠메푸드.svg",             # 쿠메푸드
        # row 4: 추가 4개 (다양성 grid)
        "cupora.png",              # 쿠포라 (무한건축 negative logo 대체)
        "스튜디오에디.svg",          # 스튜디오에디
        "핏틀리.svg",              # 핏틀리
        "뉴레드.svg",              # 뉴레드
    ]

    # ── 4×4 grid 좌표
    grid_left = 0.6
    grid_top = 1.75
    grid_w = 13.333 - 2 * grid_left
    grid_h = 4.85  # bottom 6.6 — footer divider 6.72 위 여유
    col_w = grid_w / 4
    row_h = grid_h / 4

    logo_dir = ASSETS / "case-logo"
    for i, fname in enumerate(partners):
        row = i // 4
        col = i % 4
        cell_left = grid_left + col * col_w
        cell_top = grid_top + row * row_h
        src = logo_dir / fname
        if src.exists():
            white = _ensure_white_logo(src)
            _fit_logo(slide, white, cell_left, cell_top, col_w, row_h)


def build_sox(prs, slide_no, label_en, label_ko, subtitle, examples):
    """s7·s8·s9 공통 — SoR/SoE/SoI 패턴.

    위계: 영문 약어 punch(흰, 메인) → 한글 단어(muted) → 부제 1줄(흰, 풍부) → 예시(muted, 작게).
    """
    slide = add_black_slide(prs)
    add_brandlogy(slide, slide_no)

    add_punch_text(slide, label_en, 200, top_inches=1.6, height_inches=2.6)
    add_caption(slide, label_ko, top_inches=4.45, font_size_pt=44, color=MUTED_GRAY)
    add_caption(slide, subtitle, top_inches=5.15, font_size_pt=22, color=INK_WHITE)
    add_caption(slide, examples, top_inches=5.85, font_size_pt=18)


def build_pyramid_diagram(prs, slide_no=7):
    """피라미드 다이어그램 — SoR(토대)·SoE(중간)·SoI(정점) 3-layer.

    좌측: SVG 피라미드 (homepage CoreSystemGraphic·LayerSection 노하우 풀패스
          — currentColor + stroke 위계 + fillOpacity 깊이 + dim line + cut mark + 라벤더 정점).
    우측: 3-row 라벨 컬럼 (eyebrow LAYER 0X / 영문 punch + 한글).
    배치: s7 (DX·AX 다음, SoR/SoE/SoI 정의 슬라이드 앞 — 전체 그림 먼저).
    """
    _ensure_s10_png()
    slide = add_black_slide(prs)
    add_brandlogy(slide, slide_no)

    # ── 좌측: 피라미드 PNG (5x5", SVG viewBox 600x600)
    pyr_left = Inches(0.7)
    pyr_top = Inches(0.95)
    pyr_w = Inches(5.0)
    pyr_h = Inches(5.0)
    if S10_PYRAMID_PNG.exists():
        slide.shapes.add_picture(
            str(S10_PYRAMID_PNG), pyr_left, pyr_top, width=pyr_w, height=pyr_h
        )

    # ── 우측: 라벨 컬럼 — 각 피라미드 layer 영역 top과 align (stack from top)
    # SVG y (60~540) → PPT y (pyr_top + svg_y * 5/600 inches)
    unit_in = 5.0 / 600  # inches per SVG unit
    layers = [
        # (eyebrow, 영문 punch, 한글, svg_layer_center_y)
        # 부제는 화면 텍스트 최소화 원칙(v10 spec)에 따라 제거 — 발화로 풀음
        ("LAYER 03", "SoI", "지능", 140),
        ("LAYER 02", "SoE", "참여", 300),
        ("LAYER 01", "SoR", "기록", 460),
    ]

    label_left = Inches(6.5)
    label_w = Inches(6.3)
    row_h = Inches(1.1)

    for eyebrow, label_en, label_ko, svg_center in layers:
        # 각 layer 영역 center에 라벨 vertical center align
        center_in = 0.95 + svg_center * unit_in
        tb = slide.shapes.add_textbox(
            label_left, Inches(center_in - 0.55), label_w, row_h
        )
        tf = tb.text_frame
        tf.margin_left = tf.margin_right = 0
        tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True

        # Line 1: eyebrow (LAYER 0X) — mono small caps
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.LEFT
        r1 = p1.add_run()
        r1.text = eyebrow
        set_font(r1, FONT_MONO)
        r1.font.size = Pt(11)
        r1.font.color.rgb = INK_TERTIARY
        set_letter_spacing(r1, 22.0)

        # Line 2: 영문 punch + 한글 inline
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        p2.space_before = Pt(6)
        r2a = p2.add_run()
        r2a.text = label_en
        set_font(r2a, FONT_BLACK)
        r2a.font.size = Pt(48)
        r2a.font.color.rgb = INK_WHITE
        set_letter_spacing(r2a, PUNCH_SPACING_PCT)

        r2b = p2.add_run()
        r2b.text = f"   {label_ko}"
        set_font(r2b, FONT_LIGHT)
        r2b.font.size = Pt(24)
        r2b.font.color.rgb = MUTED_GRAY


def build_s11_plus_ai(prs, slide_no=11):
    """+ AI — 좌측 피라미드 silhouette + 우측 "그 위에 + AI" mega punch.

    이전 슬라이드에서 본 피라미드 구조 위에 AI가 올라간다는 메타포를 한 슬라이드에 시각화.
    피라미드 SVG 재사용 (라벨은 우측 컬럼 없이 silhouette만 전달).
    """
    _ensure_s10_png()
    slide = add_black_slide(prs)
    add_brandlogy(slide, slide_no)

    # ── 좌측: 피라미드 silhouette (작게, 시각 echo)
    pyr_w = Inches(4.0)
    pyr_h = Inches(4.0)
    pyr_left = Inches(0.7)
    pyr_top = Inches(1.5)  # 콘텐츠 영역 중앙 (1.5 ~ 5.5)
    if S10_PYRAMID_PNG.exists():
        slide.shapes.add_picture(
            str(S10_PYRAMID_PNG), pyr_left, pyr_top, width=pyr_w, height=pyr_h
        )

    # ── 우측: "그 위에" eyebrow + "+ AI" mega punch
    right_left = Inches(5.2)
    right_w = CANVAS_W - right_left - Inches(0.5)

    # eyebrow "그 위에"
    eb = slide.shapes.add_textbox(right_left, Inches(2.0), right_w, Inches(0.6))
    ep = eb.text_frame.paragraphs[0]
    ep.alignment = PP_ALIGN.CENTER
    er = ep.add_run()
    er.text = "그 위에"
    set_font(er, FONT_LIGHT)
    er.font.size = Pt(32)
    er.font.color.rgb = MUTED_GRAY

    # mega punch "+ AI"
    punch_tb = slide.shapes.add_textbox(right_left, Inches(2.85), right_w, Inches(3.3))
    tf = punch_tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "+ AI"
    set_font(r, FONT_BLACK)
    r.font.size = Pt(220)
    r.font.color.rgb = INK_WHITE
    set_letter_spacing(r, PUNCH_SPACING_PCT)


def build_s12_self_audit(prs):
    """s12 우리 회사는? + 체크박스 SoR/SoE/SoI (자기진단 cliff hanger)."""
    slide = add_black_slide(prs)
    add_punch_text(slide, "우리 회사는?", 96, top_inches=1.5, height_inches=1.3)

    tb = slide.shapes.add_textbox(Inches(0), Inches(3.7), CANVAS_W, Inches(3.0))
    tf = tb.text_frame
    tf.vertical_anchor = MSO_ANCHOR.TOP

    items = ["☐  SoR — 기록", "☐  SoE — 참여", "☐  SoI — 지능"]
    for idx, text in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(16)
        r = p.add_run()
        r.text = text
        set_font(r, FONT_LIGHT)
        r.font.size = Pt(36)
        r.font.color.rgb = MUTED_GRAY


def build_s13_close(prs):
    """s13 검정 마무리 — 중앙 SSOD 가로형 워드마크 + 김지명."""
    slide = add_black_slide(prs)

    # 중앙 가로형 워드마크 (closing brand mark)
    if SSOD_LOGO.exists():
        h = Inches(1.4)
        w = Inches(1.4 * 2440 / 1585)
        left = (CANVAS_W - w) / 2
        top = (CANVAS_H - h) / 2 - Inches(0.3)
        slide.shapes.add_picture(str(SSOD_LOGO), left, top, width=w, height=h)

    # 워드마크 아래 김지명 (중앙 정렬)
    tb = slide.shapes.add_textbox(Inches(0), Inches(5.1), CANVAS_W, Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "김지명"
    set_font(r, FONT_LIGHT)
    r.font.size = Pt(18)
    r.font.color.rgb = INK_WHITE


def main():
    prs = Presentation()
    prs.slide_width = CANVAS_W
    prs.slide_height = CANVAS_H

    build_s1_cover(prs)
    build_s2_promise(prs, slide_no=2)
    build_s3_years(prs, slide_no=3)
    build_s4_dual_ambassador(prs, slide_no=4)
    build_s5_scale(prs, slide_no=5)
    build_brand_wall(prs, slide_no=6)
    build_s6_dx_ax(prs, slide_no=7)
    build_pyramid_diagram(prs, slide_no=8)
    build_sox(prs, 9, "SoR", "기록", "회사의 거래·사실이 권위 있게 적힌 자리", "ERP · CRM · 매출")
    build_sox(prs, 10, "SoE", "참여", "사람의 상호작용·협업이 일어나는 자리", "카톡 · 노션 · 메일")
    build_sox(prs, 11, "SoI", "지능", "여러 데이터를 합쳐 데이터를 만드는 자리", "시니어의 노하우 · 사장님의 영업 비결 등")
    build_s11_plus_ai(prs, slide_no=12)
    build_s12_self_audit(prs)
    build_s13_close(prs)

    prs.save(OUTPUT)
    print(f"saved → {OUTPUT}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
